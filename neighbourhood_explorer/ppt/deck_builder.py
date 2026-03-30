from __future__ import annotations

from collections import defaultdict
from collections.abc import Iterable
from dataclasses import dataclass
from io import BytesIO
import logging

import pandas as pd
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from neighbourhood_explorer.data_access import (
    available_periods,
    comparison_bundle,
    composition_context_frame,
    indicator_frame,
    indicator_metadata,
    indicator_timeseries_bundle,
    load_neighbourhood_reference,
    map_value_frame,
    ranked_distribution,
)
from neighbourhood_explorer.footprints import resolve_footprint_label

from .chart_registry import normalise_ppt_visual_name, render_primary_visual
from .indicator_layouts import (
    IndicatorSlideContent,
    LayoutPlan,
    choose_indicator_layout,
    validate_slide_layout,
)
from .layout_engine import (
    add_badge,
    add_bullets,
    add_fitted_textbox,
    add_footer_rule,
    add_logo,
    add_metric_card,
    add_page_number,
    add_panel,
    add_textbox,
    apply_slide_background,
    create_presentation,
)
from .style_tokens import THEME
from .text_summariser import (
    clean_text,
    comparison_summary_text,
    contents_examples,
    compress_note_text,
    full_footer_method_note,
    format_indicator_value,
    merge_narrative_text,
    neighbourhood_name_summary,
    percentile_note,
    pluralise,
    recent_change_note,
    report_date_label,
    selection_label,
    selection_scope_summary,
    short_aggregation_label,
    short_caveat,
    short_source_label,
    shorten_display_title,
    truncate_text,
)


logger = logging.getLogger(__name__)


@dataclass
class ReportContext:
    selected_ids: list[str]
    indicator_ids: list[str]
    current_indicator_id: str
    current_period: str
    include_borough: bool
    include_london: bool
    include_overview_slide: bool
    include_detail_tables: bool
    include_methodology_slide: bool
    report_title: str
    configured_topics: list[str]
    selected_meta: pd.DataFrame
    selection_names: str
    period_map: dict[str, str]
    indicator_visual_overrides: dict[str, str]
    indicator_visual_config_overrides: dict[str, dict[str, str]]


@dataclass
class MetricSlideRenderState:
    plan: LayoutPlan
    text_results: dict[str, object]
    visual_result: object
    validation: object
    primary_text: str
    secondary_text: str
    merged_text: str


def build_powerpoint_report(
    *,
    selected_ids: list[str] | Iterable[str],
    indicator_ids: list[str],
    current_indicator_id: str,
    current_period: str,
    include_borough: bool,
    include_london: bool,
    include_overview_slide: bool = True,
    include_detail_tables: bool = True,
    include_methodology_slide: bool = True,
    report_title: str | None = None,
    configured_topics: list[str] | None = None,
    indicator_visual_overrides: dict[str, str] | None = None,
    indicator_visual_config_overrides: dict[str, dict[str, str]] | None = None,
) -> bytes:
    selected_list = [str(item) for item in selected_ids]
    if not selected_list:
        raise ValueError("At least one neighbourhood must be selected.")
    if not indicator_ids:
        raise ValueError("At least one indicator must be selected.")

    selected_meta, selection_names = _selected_neighbourhood_summary(selected_list)
    period_map = _resolve_period_map(indicator_ids, current_indicator_id, current_period)
    title = report_title or f"London neighbourhood report - {selection_names[:70]}"
    context = ReportContext(
        selected_ids=selected_list,
        indicator_ids=indicator_ids,
        current_indicator_id=current_indicator_id,
        current_period=current_period,
        include_borough=include_borough,
        include_london=include_london,
        include_overview_slide=include_overview_slide,
        include_detail_tables=include_detail_tables,
        include_methodology_slide=include_methodology_slide,
        report_title=title,
        configured_topics=list(configured_topics or []),
        selected_meta=selected_meta,
        selection_names=selection_names,
        period_map=period_map,
        indicator_visual_overrides={str(key): str(value) for key, value in (indicator_visual_overrides or {}).items()},
        indicator_visual_config_overrides={
            str(indicator_id): {str(key): str(value) for key, value in config.items()}
            for indicator_id, config in (indicator_visual_config_overrides or {}).items()
        },
    )

    prs = create_presentation(THEME)
    slide_number = 1

    _build_cover_slide(prs, context)
    if context.include_overview_slide:
        slide_number += 1
        _build_selection_summary_slide(prs, context, slide_number)
        slide_number += 1
        _build_contents_slide(prs, context, slide_number)

    grouped = _group_indicator_metadata(context)
    for category, entries in grouped.items():
        slide_number += 1
        _build_category_divider_slide(prs, category, entries, slide_number)
        for entry in entries:
            slide_number += 1
            _build_metric_slide(prs, context, entry, slide_number)

    if context.include_methodology_slide:
        slide_number += 1
        _build_methodology_slide(prs, context, slide_number)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _build_cover_slide(prs, context: ReportContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, width=Inches(1.3))
    accent = slide.shapes.add_shape(1, THEME.margin_x, Inches(1.35), Inches(0.12), Inches(1.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME.navy
    accent.line.color.rgb = THEME.navy
    add_textbox(
        slide,
        THEME.margin_x + Inches(0.24),
        Inches(1.2),
        Inches(8.8),
        Inches(0.9),
        context.report_title,
        font_size=30 if len(context.report_title) < 72 else 26,
        bold=True,
        font_name=THEME.font_family_bold,
    )
    add_textbox(
        slide,
        THEME.margin_x + Inches(0.24),
        Inches(2.15),
        Inches(6.8),
        Inches(0.34),
        "Neighbourhood summary and benchmark report for the current London footprint",
        font_size=15,
        color=THEME.muted,
    )
    add_panel(slide, Inches(8.95), Inches(1.05), Inches(3.9), Inches(1.9), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(9.18), Inches(1.28), Inches(3.35), Inches(0.24), "Selection", font_size=11, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    add_textbox(slide, Inches(9.18), Inches(1.58), Inches(3.2), Inches(0.5), selection_scope_summary(context.selected_meta), font_size=20, bold=True, font_name=THEME.font_family_bold)
    add_textbox(slide, Inches(9.18), Inches(2.18), Inches(3.2), Inches(0.32), f"Generated {report_date_label()}", font_size=11, color=THEME.muted)
    add_panel(slide, THEME.margin_x, Inches(3.4), Inches(5.85), Inches(2.85), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(0.72), Inches(3.68), Inches(2.3), Inches(0.24), "Footprint in this deck", font_size=11, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    add_textbox(
        slide,
        Inches(0.72),
        Inches(4.0),
        Inches(5.2),
        Inches(0.46),
        selection_scope_summary(context.selected_meta),
        font_size=18,
        bold=True,
        font_name=THEME.font_family_bold,
    )
    add_textbox(
        slide,
        Inches(0.72),
        Inches(4.56),
        Inches(5.1),
        Inches(0.46),
        truncate_text(_selection_examples_text(context.selected_meta), 105),
        font_size=12,
        color=THEME.text,
    )

    categories = context.configured_topics or list(_group_indicator_metadata(context).keys())
    add_panel(slide, Inches(6.58), Inches(3.4), Inches(6.15), Inches(2.85), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(6.82), Inches(3.68), Inches(2.5), Inches(0.24), "Themes in this deck", font_size=11, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    y = Inches(4.0)
    x = Inches(6.84)
    for label in categories[:5]:
        badge_width = Inches(min(2.55, max(1.6, 0.12 * len(label))))
        add_badge(slide, x, y, label, fill=THEME.category_accent(label), color=THEME.surface, theme=THEME, width=badge_width)
        x += badge_width + Inches(0.1)
        if x > Inches(11.3):
            x = Inches(6.84)
            y += Inches(0.42)
    comparisons = []
    if context.include_borough:
        comparisons.append("borough benchmarks")
    if context.include_london:
        comparisons.append("London benchmark")
    add_textbox(
        slide,
        Inches(6.84),
        Inches(5.32),
        Inches(5.5),
        Inches(0.32),
        f"Comparisons: {', '.join(comparisons) if comparisons else 'none selected'}",
        font_size=12,
        color=THEME.text,
    )


def _build_selection_summary_slide(prs, context: ReportContext, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, right_aligned=True, width=Inches(0.9))
    add_badge(slide, THEME.margin_x, Inches(0.36), "Selection summary", fill=THEME.navy, color=THEME.surface, theme=THEME, width=Inches(1.8))
    add_textbox(slide, THEME.margin_x, Inches(0.78), Inches(8.6), Inches(0.36), "What this report covers", font_size=24, bold=True, font_name=THEME.font_family_bold)
    neighbourhood_count = int(context.selected_meta["neighbourhood_id"].nunique())
    borough_count = len({item.strip() for raw in context.selected_meta["borough_name"].astype(str) for item in raw.split(";") if item.strip()})
    card_y = Inches(1.35)
    card_width = Inches(2.95)
    card_gap = Inches(0.18)
    cards = [
        ("Neighbourhoods", str(neighbourhood_count), "In the selected footprint", THEME.navy),
        ("Boroughs", str(borough_count), "Covered by the selection", THEME.teal),
        ("Indicators", str(len(context.indicator_ids)), "Included in this deck", THEME.blue),
        ("Themes", str(len(context.configured_topics or _group_indicator_metadata(context))), "Grouped for exploration", THEME.gold),
    ]
    for idx, (label, value, caption, accent) in enumerate(cards):
        add_metric_card(
            slide,
            THEME.margin_x + (card_width + card_gap) * idx,
            card_y,
            card_width,
            Inches(0.95),
            label=label,
            value=value,
            caption=caption,
            accent=accent,
            theme=THEME,
        )
    add_panel(slide, THEME.margin_x, Inches(2.55), Inches(6.1), Inches(3.95), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(0.7), Inches(2.78), Inches(3.0), Inches(0.22), "Footprint at a glance", font_size=12, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    add_textbox(
        slide,
        Inches(0.7),
        Inches(3.14),
        Inches(5.45),
        Inches(0.42),
        selection_scope_summary(context.selected_meta),
        font_size=20,
        bold=True,
        font_name=THEME.font_family_bold,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(3.7),
        Inches(5.3),
        Inches(0.46),
        truncate_text(_selection_examples_text(context.selected_meta), 115),
        font_size=12,
        color=THEME.text,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(4.38),
        Inches(5.3),
        Inches(0.44),
        f"Comparisons included: {_comparison_label(context)}.",
        font_size=12,
        color=THEME.text,
    )

    add_panel(slide, Inches(6.8), Inches(2.55), Inches(6.05), Inches(3.95), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(7.03), Inches(2.78), Inches(2.8), Inches(0.22), "How the deck is organised", font_size=12, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    bullets = [
        "Each theme is grouped into a clear section with one slide per metric.",
        "Selected footprint values are shown first, with London and borough benchmarks where enabled.",
        "Monthly crime indicators prioritise trend; snapshot indicators prioritise comparison.",
        "Methodology notes are kept short on each slide, with a fuller summary at the end.",
    ]
    add_bullets(slide, Inches(7.0), Inches(3.1), Inches(5.35), Inches(1.75), bullets, font_size=11, theme=THEME, color=THEME.text)
    add_textbox(
        slide,
        Inches(7.03),
        Inches(5.18),
        Inches(5.35),
        Inches(0.44),
        f"Comparisons enabled: {_comparison_label(context)}",
        font_size=12,
        bold=True,
        font_name=THEME.font_family_bold,
    )
    add_page_number(slide, slide_number, THEME)


def _build_contents_slide(prs, context: ReportContext, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, right_aligned=True, width=Inches(0.9))
    add_badge(slide, THEME.margin_x, Inches(0.36), "Report contents", fill=THEME.navy, color=THEME.surface, theme=THEME, width=Inches(1.8))
    add_textbox(slide, THEME.margin_x, Inches(0.78), Inches(8.4), Inches(0.36), "Included themes and example indicators", font_size=24, bold=True, font_name=THEME.font_family_bold)
    grouped = _group_indicator_metadata(context)
    card_width = Inches(5.92)
    card_height = Inches(1.62)
    column_positions = [THEME.margin_x, Inches(6.72)]
    row_tops = [Inches(1.45), Inches(3.33), Inches(5.21)]
    positions = [(left, top) for top in row_tops for left in column_positions]
    for (category, entries), (left, top) in zip(grouped.items(), positions, strict=False):
        accent = THEME.category_accent(category)
        add_panel(slide, left, top, card_width, card_height, theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_badge(slide, left + Inches(0.18), top + Inches(0.18), category, fill=accent, color=THEME.surface, theme=THEME)
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.66),
            Inches(1.4),
            Inches(0.3),
            pluralise(len(entries), "indicator"),
            font_size=18,
            bold=True,
            font_name=THEME.font_family_bold,
        )
        titles = [str(entry["title"]) for entry in entries]
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.08),
            Inches(5.35),
            Inches(0.3),
            contents_examples(titles, max_examples=4),
            font_size=11,
            color=THEME.text,
        )
        modules = sorted({clean_text(entry.get("module_label")) for entry in entries if clean_text(entry.get("module_label"))})
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.33),
            Inches(5.3),
            Inches(0.22),
            f"Modules: {', '.join(modules[:3])}" + ("..." if len(modules) > 3 else ""),
            font_size=10,
            color=THEME.muted,
        )
    add_page_number(slide, slide_number, THEME)


def _build_category_divider_slide(prs, category: str, entries: list[dict[str, object]], slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    accent = THEME.category_accent(category)
    block = slide.shapes.add_shape(1, THEME.margin_x, Inches(1.15), Inches(0.16), Inches(4.95))
    block.fill.solid()
    block.fill.fore_color.rgb = accent
    block.line.color.rgb = accent
    add_textbox(
        slide,
        Inches(1.0),
        Inches(1.45),
        Inches(7.0),
        Inches(0.42),
        category,
        font_size=28,
        bold=True,
        font_name=THEME.font_family_bold,
    )
    modules = sorted({clean_text(entry.get("module_label")) for entry in entries if clean_text(entry.get("module_label"))})
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.02),
        Inches(8.0),
        Inches(0.34),
        f"{pluralise(len(entries), 'indicator')} across {pluralise(len(modules), 'module')}",
        font_size=16,
        color=THEME.muted,
    )
    add_panel(slide, Inches(8.8), Inches(1.28), Inches(3.9), Inches(1.65), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(9.05), Inches(1.56), Inches(3.2), Inches(0.24), "Included modules", font_size=11, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    add_bullets(slide, Inches(9.0), Inches(1.9), Inches(3.25), Inches(0.9), modules[:4], font_size=10, theme=THEME, color=THEME.text)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(3.15),
        Inches(10.8),
        Inches(0.46),
        truncate_text(_category_explainer(category), 170),
        font_size=16,
        color=THEME.text,
    )
    add_page_number(slide, slide_number, THEME)


def _build_metric_slide(prs, context: ReportContext, meta: dict[str, object], slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, right_aligned=True, width=Inches(0.86))

    indicator_id = str(meta["indicator_id"])
    period = context.period_map[indicator_id]
    selected_visual = _normalise_primary_visual(
        context.indicator_visual_overrides.get(indicator_id) or meta.get("primary_visualisation") or "benchmark_lollipop"
    )
    selected_visual_config = context.indicator_visual_config_overrides.get(indicator_id, {})
    bundle = comparison_bundle(indicator_id, period, context.selected_ids)
    selection = bundle["selection"]
    borough_df = bundle["borough_benchmarks"] if context.include_borough else pd.DataFrame()
    london_df = bundle["london_benchmark"] if context.include_london else pd.DataFrame()
    comparison_rows, benchmark_rows = _comparison_rows(selection, borough_df, london_df, context.selected_ids)
    detail_rows = _selected_detail_rows(indicator_id, period, context.selected_ids)
    composition_df = (
        composition_context_frame(
            indicator_id,
            period,
            context.selected_ids,
            include_borough=context.include_borough,
            include_london=context.include_london,
        )
        if selected_visual == "population_pyramid"
        else pd.DataFrame()
    )
    timeseries = (
        indicator_timeseries_bundle(
            indicator_id,
            context.selected_ids,
            include_borough=context.include_borough,
            include_london=context.include_london,
        )
        if selected_visual == "monthly_trend"
        else pd.DataFrame()
    )
    should_load_ranked = (
        str(meta.get("secondary_visualisation")) == "ranked_distribution"
        or selected_visual == "distribution_band"
        or selected_visual == "context_strip"
        or (
            selected_visual not in {"choropleth_map", "monthly_trend"}
            and benchmark_rows.empty
        )
    )
    ranked = ranked_distribution(indicator_id, period) if should_load_ranked else pd.DataFrame()
    map_values = map_value_frame(indicator_id, period) if selected_visual == "choropleth_map" else pd.DataFrame()

    accent = THEME.category_accent(clean_text(meta.get("top_level_category")))
    cards = _summary_cards(selection, benchmark_rows, context, meta) if selection is not None else []
    subtitle = f"{clean_text(meta.get('module_label'))} | {short_source_label(meta, period)}"

    if selection is None:
        logger.info(
            "PPT slide %s:%s layout=no_data visual=%s title_len=%s subtitle_len=%s",
            indicator_id,
            period,
            selected_visual,
            len(shorten_display_title(meta.get("title"), meta.get("module_label"), meta.get("source_name"), max_chars=THEME.layout.max_display_title_chars)),
            len(subtitle),
        )
        _render_metric_header(
            slide,
            category_label=clean_text(meta.get("top_level_category")) or "Indicator",
            accent=accent,
            title=shorten_display_title(meta.get("title"), meta.get("module_label"), meta.get("source_name"), max_chars=THEME.layout.max_display_title_chars),
            subtitle=subtitle,
            text_results={},
            title_variant="narrative_heavy",
        )
        add_panel(slide, THEME.margin_x, Inches(2.45), Inches(12.25), Inches(3.65), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_fitted_textbox(
            slide,
            Inches(0.76),
            Inches(3.22),
            Inches(11.2),
            Inches(0.48),
            "No data are available for this indicator for the current footprint and period.",
            min_font_size=15,
            max_font_size=20,
            max_lines=2,
            allow_truncation=False,
            bold=True,
            font_name=THEME.font_family_bold,
            theme=THEME,
        )
        add_fitted_textbox(
            slide,
            Inches(0.76),
            Inches(3.86),
            Inches(10.8),
            Inches(0.72),
            compress_note_text(full_footer_method_note(meta), max_chars=THEME.layout.max_footer_note_chars),
            min_font_size=10,
            max_font_size=12,
            max_lines=4,
            allow_truncation=True,
            color=THEME.muted,
            theme=THEME,
        )
        add_footer_rule(slide, THEME)
        _write_slide_notes(
            slide,
            full_title=clean_text(meta.get("title")),
            primary_note="No data were available for the selected footprint and period.",
            secondary_note="",
            footer_note=full_footer_method_note(meta),
            layout_variant="no_data",
            warnings=[],
        )
        add_page_number(slide, slide_number, THEME)
        return

    visual_type, visual_is_weak = _metric_visual_type(
        selected_visual=selected_visual,
        comparison_rows=comparison_rows,
        timeseries=timeseries,
        map_values=map_values,
        ranked=ranked,
    )
    has_borough_range = bool((benchmark_rows["kind"] == "borough").sum() > 1 if not benchmark_rows.empty else False)
    primary_note = (
        recent_change_note(timeseries, selection["unit"])
        if visual_type == "trend"
        else comparison_summary_text(selection["value"], benchmark_rows, selection["unit"], has_borough_range=has_borough_range)
    )
    predicted_visual_note = _predicted_visual_note(
        selected_visual=selected_visual,
        selection=selection,
        benchmark_rows=benchmark_rows,
        ranked=ranked,
        timeseries=timeseries,
        map_values=map_values,
    )
    secondary_note = _secondary_metric_note(
        meta=meta,
        selected_visual=selected_visual,
        selection=selection,
        benchmark_rows=benchmark_rows,
        ranked=ranked,
        detail_rows=detail_rows,
        include_detail_tables=context.include_detail_tables,
        visual_note=predicted_visual_note,
    )

    content = _build_indicator_slide_content(
        meta=meta,
        indicator_id=indicator_id,
        period=period,
        subtitle=subtitle,
        primary_note=primary_note,
        secondary_note=secondary_note,
        cards=cards,
        visual_type=visual_type,
        visual_is_weak=visual_is_weak,
        has_borough_benchmark=not borough_df.empty,
        has_london_benchmark=not london_df.empty,
    )
    render_state = _render_metric_slide_pass(
        slide,
        context=context,
        meta=meta,
        content=content,
        cards=cards,
        selected_visual=selected_visual,
        selected_visual_config=selected_visual_config,
        selection=selection,
        comparison_rows=comparison_rows,
        benchmark_rows=benchmark_rows,
        timeseries=timeseries,
        ranked=ranked,
        map_values=map_values,
        composition_df=composition_df,
        detail_rows=detail_rows,
        accent=accent,
    )
    autocorrect_variant = _choose_metric_autocorrect_variant(content, render_state.plan, render_state.validation)
    if autocorrect_variant and autocorrect_variant != render_state.plan.variant:
        logger.info(
            "PPT slide %s autocorrect rerender from=%s to=%s warnings=%s",
            content.slide_id,
            render_state.plan.variant,
            autocorrect_variant,
            " | ".join(render_state.validation.warnings),
        )
        _clear_slide_shapes(slide)
        render_state = _render_metric_slide_pass(
            slide,
            context=context,
            meta=meta,
            content=content,
            cards=cards,
            selected_visual=selected_visual,
            selected_visual_config=selected_visual_config,
            selection=selection,
            comparison_rows=comparison_rows,
            benchmark_rows=benchmark_rows,
            timeseries=timeseries,
            ranked=ranked,
            map_values=map_values,
            composition_df=composition_df,
            detail_rows=detail_rows,
            accent=accent,
            force_variant=autocorrect_variant,
        )
    if render_state.validation.warnings:
        logger.warning("PPT slide %s validation warnings: %s", content.slide_id, " | ".join(render_state.validation.warnings))
    _write_slide_notes(
        slide,
        full_title=content.full_title,
        primary_note=primary_note,
        secondary_note=render_state.secondary_text,
        footer_note=full_footer_method_note(meta),
        layout_variant=render_state.plan.variant,
        warnings=render_state.validation.warnings,
        visual_debug=getattr(render_state.visual_result, "debug", {}),
    )
    add_page_number(slide, slide_number, THEME)


def _render_metric_slide_pass(
    slide,
    *,
    context: ReportContext,
    meta: dict[str, object],
    content: IndicatorSlideContent,
    cards: list[dict[str, object]],
    selected_visual: str,
    selected_visual_config: dict[str, str],
    selection: dict[str, object],
    comparison_rows: pd.DataFrame,
    benchmark_rows: pd.DataFrame,
    timeseries: pd.DataFrame,
    ranked: pd.DataFrame,
    map_values: pd.DataFrame,
    composition_df: pd.DataFrame,
    detail_rows: pd.DataFrame,
    accent: str,
    force_variant: str | None = None,
) -> MetricSlideRenderState:
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, right_aligned=True, width=Inches(0.86))

    plan = choose_indicator_layout(content, THEME, force_variant=force_variant)
    logger.info(
        "PPT slide %s layout=%s visual=%s title_len=%s subtitle_len=%s note_len=%s/%s kpis=%s weak_visual=%s forced_variant=%s",
        content.slide_id,
        plan.variant,
        content.visual_type,
        content.title_length,
        content.subtitle_length,
        content.primary_length,
        content.secondary_length,
        content.kpi_count,
        content.visual_is_weak,
        force_variant or "",
    )

    text_results: dict[str, object] = {}
    _render_metric_header(
        slide,
        category_label=clean_text(meta.get("top_level_category")) or "Indicator",
        accent=accent,
        title=content.display_title,
        subtitle=content.subtitle,
        text_results=text_results,
        title_variant=plan.variant,
        plan=plan,
    )
    _add_metric_cards_row(slide, cards, boxes=plan.kpi_boxes, text_results=text_results)

    visual_result = render_primary_visual(
        slide,
        primary_visual=selected_visual,
        visual_config=selected_visual_config,
        selection=selection,
        comparison_rows=comparison_rows,
        benchmark_rows=benchmark_rows,
        timeseries=timeseries,
        ranked=ranked,
        map_values=map_values,
        composition_df=composition_df,
        current_ids=context.selected_ids,
        unit=selection["unit"],
        left=plan.visual_box.left,
        top=plan.visual_box.top,
        width=plan.visual_box.width,
        height=plan.visual_box.height,
        theme=THEME,
    )

    primary_text = content.primary_note
    secondary_text = compress_note_text(
        _secondary_metric_note(
            meta=meta,
            selected_visual=selected_visual,
            selection=selection,
            benchmark_rows=benchmark_rows,
            ranked=ranked,
            detail_rows=detail_rows,
            include_detail_tables=context.include_detail_tables,
            visual_note=visual_result.note,
        ),
        max_chars=THEME.layout.max_note_chars,
    )
    if visual_result.suppressed_benchmarks:
        secondary_text = compress_note_text(
            f"{secondary_text} Plus {visual_result.suppressed_benchmarks} more borough benchmarks are suppressed here for readability.",
            max_chars=THEME.layout.max_note_chars,
        )
    merged_text = merge_narrative_text(primary_text, secondary_text, max_chars=320)

    _render_metric_narratives(
        slide,
        plan=plan,
        primary_text=primary_text,
        secondary_text=secondary_text,
        merged_text=merged_text,
        text_results=text_results,
    )
    _render_metric_footer(
        slide,
        plan=plan,
        footer_note=content.footer_note,
        footer_meta=content.footer_meta,
        text_results=text_results,
    )

    validation = validate_slide_layout(
        content,
        plan,
        text_results=text_results,
        visual_placement=visual_result,
        theme=THEME,
    )
    logger.info(
        "PPT slide %s fit title=%spt/%sl/%s subtitle=%spt/%sl/%s footer=%spt/%sl/%s visual_fallback=%s visual_fill=%s",
        content.slide_id,
        getattr(text_results.get("title"), "font_size_pt", None),
        getattr(text_results.get("title"), "line_count", None),
        getattr(text_results.get("title"), "truncated", None),
        getattr(text_results.get("subtitle"), "font_size_pt", None),
        getattr(text_results.get("subtitle"), "line_count", None),
        getattr(text_results.get("subtitle"), "truncated", None),
        getattr(text_results.get("footer_note"), "font_size_pt", None),
        getattr(text_results.get("footer_note"), "line_count", None),
        getattr(text_results.get("footer_note"), "truncated", None),
        getattr(visual_result, "fallback_used", ""),
        f"{getattr(visual_result, 'rendered_content_ratio', 0):.2f}" if getattr(visual_result, "rendered_content_ratio", None) is not None else "n/a",
    )
    logger.info(
        "PPT slide %s plan_debug=%s visual_debug=%s",
        content.slide_id,
        plan.debug,
        getattr(visual_result, "debug", {}),
    )
    return MetricSlideRenderState(
        plan=plan,
        text_results=text_results,
        visual_result=visual_result,
        validation=validation,
        primary_text=primary_text,
        secondary_text=secondary_text,
        merged_text=merged_text,
    )


def _choose_metric_autocorrect_variant(
    content: IndicatorSlideContent,
    plan: LayoutPlan,
    validation,
) -> str | None:
    if not validation.warnings:
        return None
    fixable_fullness = any("Slide fullness score is low" in warning for warning in validation.warnings)
    fixable_text_fit = any(
        fragment in warning
        for warning in validation.warnings
        for fragment in (
            "estimated to overflow",
            "fell below the minimum readable size",
            "Title exceeds the maximum line count.",
        )
    )
    weak_visual_mismatch = any("Weak visual content did not trigger" in warning for warning in validation.warnings)

    if content.visual_type == "comparison" and plan.variant == "comparison" and fixable_text_fit:
        return "narrative_heavy"
    if content.visual_type == "trend" and plan.variant not in {"trend", "narrative_heavy"}:
        return "trend"
    if content.visual_type == "comparison" and plan.variant != "narrative_heavy" and weak_visual_mismatch:
        return "narrative_heavy"
    if content.visual_type == "map" and plan.variant != "map":
        return "map"
    return None


def _clear_slide_shapes(slide) -> None:
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def _build_methodology_slide(prs, context: ReportContext, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide, THEME)
    add_logo(slide, THEME, right_aligned=True, width=Inches(0.9))
    add_badge(slide, THEME.margin_x, Inches(0.36), "Methodology", fill=THEME.navy, color=THEME.surface, theme=THEME, width=Inches(1.5))
    add_textbox(slide, THEME.margin_x, Inches(0.78), Inches(8.8), Inches(0.36), "How to read this report", font_size=24, bold=True, font_name=THEME.font_family_bold)
    cards = [
        (
            "Neighbourhood geography",
            [
                "All neighbourhood values use the app's static LSOA 2021 to neighbourhood crosswalk.",
                "Selected footprints can contain one or many neighbourhoods.",
            ],
        ),
        (
            "Aggregation rules",
            [
                "Counts are summed.",
                "Shares and rates are recomputed from summed components where possible.",
                "Display-only metrics are shown as descriptive summaries rather than additive values.",
            ],
        ),
        (
            "Benchmarks",
            [
                "London and borough benchmarks are optional and follow the current app settings.",
                "Cross-borough selections do not create an invented single borough benchmark.",
            ],
        ),
    ]
    x_positions = [THEME.margin_x, Inches(4.45), Inches(8.42)]
    for (title, bullets), left in zip(cards, x_positions, strict=False):
        add_panel(slide, left, Inches(1.42), Inches(3.7), Inches(2.2), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_textbox(slide, left + Inches(0.18), Inches(1.66), Inches(3.1), Inches(0.24), title, font_size=13, bold=True, font_name=THEME.font_family_bold)
        add_bullets(slide, left + Inches(0.16), Inches(2.0), Inches(3.28), Inches(1.35), bullets, font_size=10, theme=THEME, color=THEME.text)

    source_summary = _source_summary_rows(context)
    add_panel(slide, THEME.margin_x, Inches(4.02), Inches(12.25), Inches(2.35), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(0.7), Inches(4.26), Inches(3.2), Inches(0.22), "Source coverage in this deck", font_size=12, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    add_bullets(slide, Inches(0.68), Inches(4.56), Inches(5.8), Inches(1.38), source_summary, font_size=11, theme=THEME, color=THEME.text)
    add_panel(slide, Inches(6.95), Inches(4.02), Inches(5.78), Inches(2.35), theme=THEME, fill=THEME.surface, line=THEME.soft_line)
    add_textbox(slide, Inches(7.2), Inches(4.26), Inches(3.2), Inches(0.22), "On-slide note style", font_size=12, bold=True, color=THEME.muted, font_name=THEME.font_family_bold)
    note_examples = []
    for indicator_id in context.indicator_ids[:4]:
        meta = indicator_metadata(indicator_id)
        note_examples.append(f"{truncate_text(meta['title'], 36)}: {short_aggregation_label(meta.get('aggregation_method'))}")
    add_bullets(slide, Inches(7.18), Inches(4.56), Inches(5.1), Inches(1.38), note_examples, font_size=10, theme=THEME, color=THEME.text)
    add_textbox(
        slide,
        THEME.margin_x,
        Inches(6.58),
        Inches(10.8),
        Inches(0.22),
        "Longer source caveats stay in the app and methodology surface so each slide remains readable.",
        font_size=9,
        color=THEME.muted,
    )
    add_page_number(slide, slide_number, THEME)


def _selected_neighbourhood_summary(selected_ids: list[str]) -> tuple[pd.DataFrame, str]:
    reference = load_neighbourhood_reference()
    selected = reference[reference["neighbourhood_id"].astype(str).isin({str(item) for item in selected_ids})].copy()
    selected = selected.sort_values(["borough_name", "neighbourhood_name"]).reset_index(drop=True)
    label = resolve_footprint_label(selected_ids, reference).label
    return selected, label


def _resolve_period_map(indicator_ids: list[str], current_indicator_id: str, current_period: str) -> dict[str, str]:
    period_map: dict[str, str] = {}
    for indicator_id in indicator_ids:
        periods = available_periods(indicator_id)
        if not periods:
            period_map[indicator_id] = ""
            continue
        if indicator_id == current_indicator_id and current_period in periods:
            period_map[indicator_id] = current_period
        else:
            period_map[indicator_id] = periods[-1]
    return period_map


def _group_indicator_metadata(context: ReportContext) -> dict[str, list[dict[str, object]]]:
    grouped: defaultdict[str, list[dict[str, object]]] = defaultdict(list)
    ordered = []
    for indicator_id in context.indicator_ids:
        meta = indicator_metadata(indicator_id)
        meta["period"] = context.period_map[indicator_id]
        ordered.append(meta)
    ordered.sort(
        key=lambda item: (
            int(item.get("category_sort_order", 999)),
            int(item.get("module_sort_order", 999)),
            clean_text(item.get("title")),
        )
    )
    for meta in ordered:
        grouped[clean_text(meta.get("top_level_category")) or "Other"].append(meta)
    return dict(grouped)


def _comparison_rows(
    selection: dict[str, object] | None,
    borough_df: pd.DataFrame,
    london_df: pd.DataFrame,
    selected_ids: list[str],
) -> tuple[pd.DataFrame, pd.DataFrame]:
    if selection is None:
        empty = pd.DataFrame(columns=["label", "value", "kind"])
        return empty, empty
    comparison_rows: list[dict[str, object]] = [
        {"label": selection_label(selected_ids), "value": selection["value"], "kind": "selection"}
    ]
    benchmark_rows: list[dict[str, object]] = []
    if not borough_df.empty:
        borough_records = [
            {"label": str(row.benchmark_name), "value": row.value, "kind": "borough"} for row in borough_df.itertuples(index=False)
        ]
        benchmark_rows.extend(borough_records)
        comparison_rows.extend(borough_records)
    if not london_df.empty:
        row = london_df.iloc[0]
        london_row = {"label": "London overall", "value": row["value"], "kind": "london"}
        benchmark_rows.append(london_row)
        comparison_rows.append(london_row)
    return pd.DataFrame(comparison_rows), pd.DataFrame(benchmark_rows)


def _selected_detail_rows(indicator_id: str, period: str, selected_ids: list[str]) -> pd.DataFrame:
    frame = indicator_frame(indicator_id, period).copy()
    selected = frame[frame["neighbourhood_id"].astype(str).isin({str(item) for item in selected_ids})].copy()
    if selected.empty:
        return selected
    selected = selected.sort_values(["value", "neighbourhood_name"], ascending=[False, True]).reset_index(drop=True)
    selected["display_value"] = [
        format_indicator_value(row.value, row.unit) for row in selected[["value", "unit"]].itertuples(index=False)
    ]
    return selected[["neighbourhood_name", "borough_name", "display_value", "value"]]


def _summary_cards(selection: dict[str, object], benchmark_rows: pd.DataFrame, context: ReportContext, meta: dict[str, object]) -> list[dict[str, object]]:
    selection_descriptor = resolve_footprint_label(context.selected_ids, load_neighbourhood_reference())
    selected_count = int(selection.get("selected_neighbourhood_count", len(context.selected_ids)) or len(context.selected_ids))
    cards = [
        {
            "label": selection_label(context.selected_ids),
            "value": format_indicator_value(selection["value"], selection["unit"]),
            "caption": _selection_card_caption(selection_descriptor.kind, selection_descriptor.exact_boundary, selected_count),
            "accent": THEME.navy,
        }
    ]
    borough_rows = benchmark_rows[benchmark_rows["kind"] == "borough"] if not benchmark_rows.empty else pd.DataFrame()
    london_rows = benchmark_rows[benchmark_rows["kind"] == "london"] if not benchmark_rows.empty else pd.DataFrame()
    if not london_rows.empty:
        cards.append(
            {
                "label": "London benchmark",
                "value": format_indicator_value(london_rows.iloc[0]["value"], selection["unit"]),
                "caption": "All London neighbourhoods combined",
                "accent": THEME.orange,
            }
        )
    if len(borough_rows) == 1:
        cards.append(
            {
                "label": "Borough benchmark",
                "value": format_indicator_value(borough_rows.iloc[0]["value"], selection["unit"]),
                "caption": clean_text(borough_rows.iloc[0]["label"]),
                "accent": THEME.teal,
            }
        )
    elif len(borough_rows) > 1:
        borough_labels = [clean_text(value) for value in borough_rows["label"].astype(str).tolist()]
        min_value = borough_rows["value"].min()
        max_value = borough_rows["value"].max()
        min_display = format_indicator_value(min_value, selection["unit"])
        max_display = format_indicator_value(max_value, selection["unit"])
        no_variation = pd.notna(min_value) and pd.notna(max_value) and (
            abs(float(max_value) - float(min_value)) <= 1e-12 or min_display == max_display
        )
        cards.append(
            {
                "label": "Borough benchmarks",
                "value": min_display if no_variation else f"{min_display} to {max_display}",
                "caption": "No variation across selected borough benchmarks" if no_variation else ", ".join(borough_labels),
                "accent": THEME.teal,
            }
        )
    return cards[:3]


def _selection_card_caption(selection_kind: str, exact_boundary: bool, selected_count: int) -> str:
    if selection_kind == "region":
        return "London-wide footprint"
    if selected_count <= 1:
        return "Single neighbourhood in scope"
    if exact_boundary and selection_kind == "place":
        return f"Exact place boundary ({selected_count} neighbourhoods)"
    if exact_boundary and selection_kind == "system":
        return f"Exact system boundary ({selected_count} neighbourhoods)"
    return f"{selected_count} neighbourhoods in scope"


def _build_indicator_slide_content(
    *,
    meta: dict[str, object],
    indicator_id: str,
    period: str,
    subtitle: str,
    primary_note: str,
    secondary_note: str,
    cards: list[dict[str, object]],
    visual_type: str,
    visual_is_weak: bool,
    has_borough_benchmark: bool,
    has_london_benchmark: bool,
) -> IndicatorSlideContent:
    display_title = shorten_display_title(
        meta.get("title"),
        meta.get("module_label"),
        meta.get("source_name"),
        max_chars=THEME.layout.max_display_title_chars,
    )
    primary_text = compress_note_text(primary_note, max_chars=THEME.layout.max_note_chars)
    secondary_text = compress_note_text(secondary_note, max_chars=THEME.layout.max_note_chars)
    if secondary_text and primary_text and secondary_text.lower() == primary_text.lower():
        secondary_text = compress_note_text(short_caveat(meta), max_chars=THEME.layout.max_note_chars)
    return IndicatorSlideContent(
        slide_id=f"{indicator_id}:{period}",
        title=clean_text(meta.get("title")),
        display_title=display_title,
        full_title=clean_text(meta.get("title")),
        subtitle=subtitle,
        primary_note=primary_text,
        secondary_note=secondary_text,
        footer_note=compress_note_text(full_footer_method_note(meta), max_chars=THEME.layout.max_footer_note_chars),
        footer_meta=short_source_label(meta, period),
        visual_type=visual_type,
        visual_is_weak=visual_is_weak,
        kpi_cards=cards,
        has_borough_benchmark=has_borough_benchmark,
        has_london_benchmark=has_london_benchmark,
        visual_aspect_ratio=1.93 if visual_type == "map" else None,
    )


def _normalise_primary_visual(selected_visual: object) -> str:
    return normalise_ppt_visual_name(selected_visual)


def _metric_visual_type(
    *,
    selected_visual: str,
    comparison_rows: pd.DataFrame,
    timeseries: pd.DataFrame,
    map_values: pd.DataFrame,
    ranked: pd.DataFrame,
) -> tuple[str, bool]:
    if selected_visual == "choropleth_map" and not map_values.empty:
        return "map", False
    if selected_visual == "monthly_trend" and not timeseries.empty:
        weak = len(timeseries["period"].drop_duplicates()) < 5
        return "trend", weak
    if selected_visual == "context_strip":
        return "comparison", ranked.empty
    weak = comparison_rows.empty and ranked.empty
    return "comparison", weak


def _predicted_visual_note(
    *,
    selected_visual: str,
    selection: dict[str, object],
    benchmark_rows: pd.DataFrame,
    ranked: pd.DataFrame,
    timeseries: pd.DataFrame,
    map_values: pd.DataFrame,
) -> str:
    if selected_visual == "monthly_trend" and not timeseries.empty:
        return "Monthly trend shown across the latest available data series."
    if selected_visual == "context_strip":
        return percentile_note(ranked, selection.get("value") if selection else None)
    if selected_visual == "distribution_band":
        return "Decile band shown as a descriptive summary only."
    if selected_visual == "choropleth_map" and not map_values.empty:
        return "Neighbourhood values are shown on the London map."
    return percentile_note(ranked, selection.get("value") if selection else None)


def _render_metric_header(
    slide,
    *,
    category_label: str,
    accent,
    title: str,
    subtitle: str,
    text_results: dict[str, object],
    title_variant: str,
    plan: LayoutPlan | None = None,
) -> None:
    add_badge(slide, THEME.margin_x, Inches(0.34), category_label or "Indicator", fill=accent, color=THEME.surface, theme=THEME)
    if plan is None:
        title_shape, title_fit = add_fitted_textbox(
            slide,
            THEME.margin_x,
            Inches(0.74),
            Inches(9.25),
            Inches(0.7 if title_variant == "narrative_heavy" else 0.54),
            title,
            min_font_size=THEME.layout.title_min_font,
            max_font_size=THEME.layout.title_max_font,
            max_lines=THEME.layout.max_title_lines,
            allow_truncation=True,
            bold=True,
            font_name=THEME.font_family_bold,
            theme=THEME,
        )
        subtitle_shape, subtitle_fit = add_fitted_textbox(
            slide,
            THEME.margin_x,
            Inches(1.26),
            Inches(10.0),
            Inches(0.34),
            subtitle,
            min_font_size=THEME.layout.caption_min_font,
            max_font_size=10,
            max_lines=THEME.layout.max_subtitle_lines,
            allow_truncation=True,
            color=THEME.muted,
            theme=THEME,
        )
    else:
        title_shape, title_fit = add_fitted_textbox(
            slide,
            plan.title_box.left,
            plan.title_box.top,
            plan.title_box.width,
            plan.title_box.height,
            title,
            min_font_size=plan.title_spec.min_font_size,
            max_font_size=plan.title_spec.max_font_size,
            max_lines=plan.title_spec.max_lines,
            allow_truncation=plan.title_spec.allow_truncation,
            bold=True,
            font_name=THEME.font_family_bold,
            theme=THEME,
        )
        subtitle_shape, subtitle_fit = add_fitted_textbox(
            slide,
            plan.subtitle_box.left,
            plan.subtitle_box.top,
            plan.subtitle_box.width,
            plan.subtitle_box.height,
            subtitle,
            min_font_size=plan.subtitle_spec.min_font_size,
            max_font_size=plan.subtitle_spec.max_font_size,
            max_lines=plan.subtitle_spec.max_lines,
            allow_truncation=plan.subtitle_spec.allow_truncation,
            color=THEME.muted,
            theme=THEME,
        )
    text_results["title"] = title_fit
    text_results["subtitle"] = subtitle_fit


def _add_metric_cards_row(
    slide,
    cards: list[dict[str, object]],
    *,
    boxes: list[object],
    text_results: dict[str, object] | None = None,
) -> None:
    for index, (card, box) in enumerate(zip(cards, boxes, strict=False)):
        fits = add_metric_card(
            slide,
            box.left,
            box.top,
            box.width,
            box.height,
            label=str(card["label"]),
            value=str(card["value"]),
            caption=str(card["caption"]),
            accent=card["accent"],
            theme=THEME,
        )
        if text_results is None:
            continue
        for part, fit in fits.items():
            if fit is not None:
                text_results[f"kpi_{index}_{part}"] = fit


def _render_metric_narratives(
    slide,
    *,
    plan: LayoutPlan,
    primary_text: str,
    secondary_text: str,
    merged_text: str,
    text_results: dict[str, object],
) -> None:
    if (
        plan.merged_narrative_box is None
        and plan.primary_narrative_box is None
        and plan.secondary_narrative_box is None
    ):
        return

    if plan.merged_narrative_box is not None and plan.merged_spec is not None:
        box = plan.merged_narrative_box
        add_panel(slide, box.left, box.top, box.width, box.height, theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.08),
            Inches(2.2),
            Inches(0.18),
            "Key insight",
            font_size=9,
            bold=True,
            color=THEME.muted,
            font_name=THEME.font_family_bold,
        )
        _, fit = add_fitted_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.26),
            box.width - Inches(0.28),
            max(Inches(0.22), box.height - Inches(0.34)),
            merged_text,
            min_font_size=plan.merged_spec.min_font_size,
            max_font_size=plan.merged_spec.max_font_size,
            max_lines=plan.merged_spec.max_lines,
            allow_truncation=False,
            color=THEME.text,
            theme=THEME,
        )
        text_results["merged_narrative"] = fit
        return

    if plan.primary_narrative_box is not None and plan.primary_spec is not None:
        box = plan.primary_narrative_box
        add_panel(slide, box.left, box.top, box.width, box.height, theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.08),
            Inches(1.8),
            Inches(0.18),
            "What stands out",
            font_size=9,
            bold=True,
            color=THEME.muted,
            font_name=THEME.font_family_bold,
        )
        _, fit = add_fitted_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.26),
            box.width - Inches(0.28),
            max(Inches(0.22), box.height - Inches(0.34)),
            primary_text,
            min_font_size=plan.primary_spec.min_font_size,
            max_font_size=plan.primary_spec.max_font_size,
            max_lines=plan.primary_spec.max_lines,
            allow_truncation=False,
            color=THEME.text,
            theme=THEME,
        )
        text_results["primary_narrative"] = fit

    if plan.secondary_narrative_box is not None and plan.secondary_spec is not None:
        box = plan.secondary_narrative_box
        add_panel(slide, box.left, box.top, box.width, box.height, theme=THEME, fill=THEME.surface, line=THEME.soft_line)
        add_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.08),
            Inches(2.2),
            Inches(0.18),
            "Additional context",
            font_size=9,
            bold=True,
            color=THEME.muted,
            font_name=THEME.font_family_bold,
        )
        _, fit = add_fitted_textbox(
            slide,
            box.left + Inches(0.14),
            box.top + Inches(0.26),
            box.width - Inches(0.28),
            max(Inches(0.22), box.height - Inches(0.34)),
            secondary_text,
            min_font_size=plan.secondary_spec.min_font_size,
            max_font_size=plan.secondary_spec.max_font_size,
            max_lines=plan.secondary_spec.max_lines,
            allow_truncation=False,
            color=THEME.text,
            theme=THEME,
        )
        text_results["secondary_narrative"] = fit


def _render_metric_footer(
    slide,
    *,
    plan: LayoutPlan,
    footer_note: str,
    footer_meta: str,
    text_results: dict[str, object],
) -> None:
    add_footer_rule(slide, THEME)
    _, footer_fit = add_fitted_textbox(
        slide,
        plan.footer_note_box.left,
        plan.footer_note_box.top,
        plan.footer_note_box.width,
        plan.footer_note_box.height,
        footer_note,
        min_font_size=plan.footer_spec.min_font_size,
        max_font_size=plan.footer_spec.max_font_size,
        max_lines=plan.footer_spec.max_lines,
        allow_truncation=plan.footer_spec.allow_truncation,
        color=THEME.muted,
        theme=THEME,
    )
    _, footer_meta_fit = add_fitted_textbox(
        slide,
        plan.footer_meta_box.left,
        plan.footer_meta_box.top,
        plan.footer_meta_box.width,
        plan.footer_meta_box.height,
        footer_meta,
        min_font_size=plan.footer_meta_spec.min_font_size,
        max_font_size=plan.footer_meta_spec.max_font_size,
        max_lines=plan.footer_meta_spec.max_lines,
        allow_truncation=plan.footer_meta_spec.allow_truncation,
        color=THEME.muted,
        align=PP_ALIGN.RIGHT,
        theme=THEME,
    )
    text_results["footer_note"] = footer_fit
    text_results["footer_meta"] = footer_meta_fit


def _write_slide_notes(
    slide,
    *,
    full_title: str,
    primary_note: str,
    secondary_note: str,
    footer_note: str,
    layout_variant: str,
    warnings: list[str],
    visual_debug: dict[str, object] | None = None,
) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    lines = [
        f"Full title: {clean_text(full_title)}",
        f"Layout: {clean_text(layout_variant)}",
    ]
    if clean_text(primary_note):
        lines.append(f"What stands out: {clean_text(primary_note)}")
    if clean_text(secondary_note):
        lines.append(f"Additional context: {clean_text(secondary_note)}")
    if clean_text(footer_note):
        lines.append(f"Method note: {clean_text(footer_note)}")
    debug_summary = _format_visual_debug_summary(visual_debug or {})
    if debug_summary:
        lines.append(f"Visual QA: {debug_summary}")
    if warnings:
        lines.append("Validation warnings: " + " | ".join(clean_text(item) for item in warnings if clean_text(item)))

    paragraph = notes_frame.paragraphs[0]
    paragraph.text = lines[0]
    for line in lines[1:]:
        notes_frame.add_paragraph().text = line


def _format_visual_debug_summary(debug: dict[str, object]) -> str:
    if not debug:
        return ""
    preferred_keys = [
        "row_count",
        "longest_label_chars",
        "visible_series_count",
        "max_series_label_chars",
        "legend_position",
        "series_label_limit",
        "period_count",
        "source_size",
    ]
    parts = []
    for key in preferred_keys:
        value = debug.get(key)
        if value is None or value == "" or value == [] or value == {}:
            continue
        parts.append(f"{key}={value}")
    return "; ".join(parts[:5])


def _secondary_metric_note(
    *,
    meta: dict[str, object],
    selected_visual: str,
    selection: dict[str, object],
    benchmark_rows: pd.DataFrame,
    ranked: pd.DataFrame,
    detail_rows: pd.DataFrame,
    include_detail_tables: bool,
    visual_note: str,
) -> str:
    if selected_visual == "monthly_trend":
        compare_note = comparison_summary_text(selection["value"], benchmark_rows, selection["unit"])
        return compare_note
    if selected_visual in {"distribution_band", "context_strip"} and visual_note:
        return visual_note
    if include_detail_tables and len(detail_rows) > 1:
        return _selection_spread_note(detail_rows)
    if visual_note:
        return visual_note
    return short_caveat(meta)


def _selection_spread_note(detail_rows: pd.DataFrame) -> str:
    if detail_rows.empty:
        return "No within-selection variation is available for this indicator."
    if len(detail_rows) == 1:
        row = detail_rows.iloc[0]
        return f"{row['neighbourhood_name']} is the only area in scope for this metric."
    highest = detail_rows.iloc[0]
    lowest = detail_rows.iloc[-1]
    if highest["neighbourhood_name"] == lowest["neighbourhood_name"]:
        return f"Selected values are aligned around {highest['display_value']}."
    return (
        f"Within the selected footprint, highest is {highest['neighbourhood_name']} ({highest['display_value']}) "
        f"and lowest is {lowest['neighbourhood_name']} ({lowest['display_value']})."
    )


def _source_summary_rows(context: ReportContext) -> list[str]:
    seen: defaultdict[str, int] = defaultdict(int)
    for indicator_id in context.indicator_ids:
        source = clean_text(indicator_metadata(indicator_id).get("source_name")) or "Public data"
        seen[source] += 1
    return [f"{source}: {count} indicators" for source, count in seen.items()]


def _category_explainer(category: str) -> str:
    explanations = {
        "Population & Demography": "Population modules show the size and profile of the selected footprint using Census 2021-based neighbourhood estimates.",
        "Health & Wellbeing": "Health modules focus on resident characteristics, disability, care needs and clinical outcomes for the selected footprint.",
        "Social Factors & Wider Determinants": "Social factors modules cover deprivation, poverty, employment, education, housing and the wider conditions that shape health and opportunity.",
        "Safety & Crime": "Crime modules prioritise monthly public police series, showing recent levels and trends for the selected footprint.",
        "Environment & Access": "Environment modules cover transport connectivity, air quality, flood risk and green infrastructure.",
    }
    return explanations.get(category, "Indicators are grouped here to keep the report easy to scan.")


def _comparison_label(context: ReportContext) -> str:
    comparisons = []
    if context.include_borough:
        comparisons.append("borough")
    if context.include_london:
        comparisons.append("London")
    return ", ".join(comparisons) if comparisons else "none"


def _selection_examples_text(selected_meta: pd.DataFrame) -> str:
    if selected_meta.empty:
        return "No neighbourhoods are currently selected."
    unique_names = selected_meta["neighbourhood_name"].dropna().astype(str).tolist()
    if len(unique_names) <= 3:
        return "Includes " + ", ".join(unique_names) + "."
    return "Includes " + neighbourhood_name_summary(selected_meta, max_items=3) + "."
