from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO

import pandas as pd
from PIL import Image, ImageChops, ImageDraw
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from neighbourhood_explorer.data_access import load_hex_icb_geography, load_map_geography
from neighbourhood_explorer.footprints import resolve_footprint_label

from .layout_engine import Box, VisualPlacement, add_fitted_textbox, add_panel, add_textbox, add_visual_contained
from .style_tokens import PptTheme, THEME
from .text_summariser import (
    chart_number_format,
    clean_text,
    format_indicator_value,
    friendly_period,
    percentile_note,
    scale_for_chart,
    truncate_text,
)


@dataclass
class VisualResult:
    note: str = ""
    suppressed_benchmarks: int = 0
    visual_kind: str = "comparison"
    aspect_ratio: float | None = None
    visual_box: Box | None = None
    visual_is_weak: bool = False
    fallback_used: str = ""
    rendered_content_ratio: float | None = None
    debug: dict[str, object] = field(default_factory=dict)


PPT_VISUAL_ALIAS_MAP = {
    "trend_line": "monthly_trend",
    "trend_line_with_rolling_average": "monthly_trend",
    "benchmark_only_compare": "benchmark_lollipop",
    "lollipop_compare": "benchmark_lollipop",
    "sorted_bar": "benchmark_lollipop",
    "stacked_100_bar": "benchmark_lollipop",
    "grouped_bar": "benchmark_lollipop",
    "crime_mix_chart": "benchmark_lollipop",
    "domain_tile_matrix": "benchmark_lollipop",
    "kpi_card": "benchmark_lollipop",
    "table_support_only": "benchmark_lollipop",
    "text_badges_or_indexed_note": "benchmark_lollipop",
    "ranked_distribution": "context_strip",
    "ranked_strip": "context_strip",
}


def normalise_ppt_visual_name(primary_visual: object) -> str:
    value = clean_text(primary_visual) or "benchmark_lollipop"
    return PPT_VISUAL_ALIAS_MAP.get(value, value)


def render_primary_visual(
    slide,
    *,
    primary_visual: str,
    visual_config: dict[str, str] | None = None,
    selection: dict[str, object] | None,
    comparison_rows: pd.DataFrame,
    benchmark_rows: pd.DataFrame,
    timeseries: pd.DataFrame,
    ranked: pd.DataFrame,
    map_values: pd.DataFrame,
    composition_df: pd.DataFrame,
    current_ids: list[str],
    unit: str | None,
    left,
    top,
    width,
    height,
    theme: PptTheme = THEME,
) -> VisualResult:
    primary_visual = normalise_ppt_visual_name(primary_visual)
    visual_config = visual_config or {}
    if primary_visual == "monthly_trend" and not timeseries.empty:
        suppressed = max(0, len(benchmark_rows[benchmark_rows["kind"] == "borough"]) - 2) if not benchmark_rows.empty else 0
        period_count = int(timeseries["period"].drop_duplicates().nunique())
        trend_debug = render_trend_chart(slide, timeseries, unit, left, top, width, height, theme=theme)
        return VisualResult(
            note="Monthly trend shown across the latest available data series.",
            suppressed_benchmarks=suppressed,
            visual_kind="trend",
            aspect_ratio=(width - Inches(0.36)) / max(height - Inches(0.6), 1),
            visual_box=Box(left, top, width, height),
            visual_is_weak=period_count < 5,
            rendered_content_ratio=0.82,
            debug={
                "period_count": period_count,
                "series_count": int(timeseries["series"].drop_duplicates().nunique()),
                **trend_debug,
            },
        )
    if primary_visual == "distribution_band":
        render_distribution_band(slide, selection, benchmark_rows, left, top, width, height, unit, theme=theme)
        return VisualResult(
            note="Decile band shown as a descriptive summary only.",
            visual_kind="comparison",
            aspect_ratio=width / max(height, 1),
            visual_box=Box(left, top, width, height),
            rendered_content_ratio=0.78,
        )
    if primary_visual == "context_strip" and not ranked.empty:
        render_context_strip(slide, ranked, selection.get("value") if selection else None, unit, left, top, width, height, theme=theme)
        return VisualResult(
            note=percentile_note(ranked, selection.get("value") if selection else None),
            visual_kind="comparison",
            aspect_ratio=width / max(height, 1),
            visual_box=Box(left, top, width, height),
            rendered_content_ratio=0.72,
            debug={"ranked_count": int(len(ranked)), "explicit_distribution_view": True},
        )
    if primary_visual == "population_pyramid" and not composition_df.empty:
        pyramid_debug = render_population_pyramid(
            slide,
            composition_df,
            comparator_subject=clean_text(visual_config.get("population_pyramid_comparator")),
            left=left,
            top=top,
            width=width,
            height=height,
            theme=theme,
        )
        if pyramid_debug.get("rendered"):
            comparator_subject = clean_text(pyramid_debug.get("comparator_subject"))
            return VisualResult(
                note=f"Population pyramid shown against {comparator_subject}." if comparator_subject else "Population pyramid shown.",
                visual_kind="comparison",
                aspect_ratio=width / max(height, 1),
                visual_box=Box(left, top, width, height),
                rendered_content_ratio=0.82,
                debug=pyramid_debug,
            )
    if primary_visual == "choropleth_map" and not map_values.empty:
        map_style = clean_text(visual_config.get("map_style")).lower() or "real"
        placement = render_choropleth_map(slide, map_values, current_ids, left, top, width, height, unit, map_style=map_style, theme=theme)
        map_target_area = max(
            (width - theme.layout.map_reserved_width) * (height - theme.layout.map_reserved_height),
            1,
        )
        return VisualResult(
            note="Neighbourhood values are shown on the London hex map." if map_style == "hex" else "Neighbourhood values are shown on the London map.",
            visual_kind="map",
            aspect_ratio=(placement.rendered_width / max(placement.rendered_height, 1)) if placement is not None else (1120 / 460),
            visual_box=placement.box if placement is not None else Box(left, top, width, height),
            rendered_content_ratio=(placement.box.area / map_target_area) if placement is not None else None,
            debug={
                "map_style": map_style,
                "source_size": [placement.source_width_px, placement.source_height_px] if placement is not None else None,
                "map_target_area": map_target_area,
            },
        )
    selection_only = (
        not comparison_rows.empty
        and comparison_rows["kind"].astype(str).fillna("").eq("selection").all()
    )
    if (comparison_rows.empty or selection_only) and not ranked.empty:
        render_context_strip(slide, ranked, selection.get("value") if selection else None, unit, left, top, width, height, theme=theme)
        return VisualResult(
            note=percentile_note(ranked, selection.get("value") if selection else None),
            visual_kind="comparison",
            aspect_ratio=width / max(height, 1),
            visual_box=Box(left, top, width, height),
            visual_is_weak=False,
            fallback_used="context_strip",
            rendered_content_ratio=0.72,
            debug={"selection_only": selection_only, "ranked_count": int(len(ranked))},
        )
    comparison_debug = render_snapshot_comparison(slide, comparison_rows, left, top, width, height, unit, theme=theme)
    return VisualResult(
        note=percentile_note(ranked, selection.get("value") if selection else None),
        visual_kind="comparison",
        aspect_ratio=width / max(height, 1),
        visual_box=Box(left, top, width, height),
        visual_is_weak=comparison_rows.empty or selection_only,
        rendered_content_ratio=0.8 if len(comparison_rows) > 1 else 0.62,
        debug={
            "row_count": int(len(comparison_rows)),
            "selection_only": selection_only,
            **comparison_debug,
        },
    )


def _normalise_visual_name(primary_visual: object) -> str:
    return normalise_ppt_visual_name(primary_visual)


def render_snapshot_comparison(slide, rows: pd.DataFrame, left, top, width, height, unit: str | None, *, theme: PptTheme) -> dict[str, object]:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.12),
        width - Inches(0.32),
        Inches(0.22),
        "How it compares",
        font_size=12,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )
    debug: dict[str, object] = {
        "row_count": int(len(rows)),
        "longest_label_chars": 0,
        "label_compacted": False,
        "label_font_max": 11,
        "row_spacing_inches": 0.0,
        "label_width_inches": 2.1,
    }
    if rows.empty:
        add_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(0.55),
            width - Inches(0.32),
            Inches(0.3),
            "No comparator values are available for this indicator.",
            font_size=12,
            color=theme.muted,
        )
        return debug

    valid = rows.copy()
    row_count = int(len(valid))
    label_lengths = [len(clean_text(getattr(row, "label"))) for row in valid.itertuples(index=False)]
    longest_label_chars = max(label_lengths, default=0)
    numeric_values: list[float] = []
    for row in valid.itertuples(index=False):
        if getattr(row, "kind") == "range":
            numeric_values.extend([float(row.value_min), float(row.value_max)])
        elif pd.notna(getattr(row, "value")):
            numeric_values.append(float(row.value))
    if not numeric_values:
        add_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(0.55),
            width - Inches(0.32),
            Inches(0.3),
            "No comparator values are available for this indicator.",
            font_size=12,
            color=theme.muted,
        )
        return debug

    axis_min = 0.0 if min(numeric_values) >= 0 else min(numeric_values) * 0.95
    axis_max = max(numeric_values) * 1.08 if max(numeric_values) else 1.0
    if axis_max <= axis_min:
        axis_max = axis_min + 1.0

    label_width = Inches(2.1)
    if longest_label_chars >= theme.layout.comparison_label_soft_chars:
        label_width = Inches(2.35)
    if longest_label_chars >= theme.layout.comparison_label_hard_chars:
        label_width = Inches(2.55)
    value_width = Inches(1.25 if row_count >= theme.layout.comparison_dense_row_threshold else 1.35)
    label_font_max = 11
    if row_count >= theme.layout.comparison_dense_row_threshold or longest_label_chars >= theme.layout.comparison_label_soft_chars:
        label_font_max = 10
    if row_count >= theme.layout.comparison_warn_row_threshold or longest_label_chars >= theme.layout.comparison_label_hard_chars:
        label_font_max = 9
    value_font_max = 10 if row_count < theme.layout.comparison_dense_row_threshold else 9
    plot_left = left + Inches(0.22) + label_width
    plot_right = left + width - Inches(0.22) - value_width
    row_top = top + Inches(0.65)
    row_spacing = min(Inches(0.72), (height - Inches(1.18)) / max(len(valid), 1))
    label_height = min(Inches(0.34), max(Inches(0.2), row_spacing - Inches(0.08)))
    value_height = min(Inches(0.32), max(Inches(0.2), row_spacing - Inches(0.1)))
    debug.update(
        {
            "row_count": row_count,
            "longest_label_chars": longest_label_chars,
            "label_compacted": label_font_max < 11 or label_width > Inches(2.1),
            "label_font_max": label_font_max,
            "row_spacing_inches": round(row_spacing / 914400, 3),
            "label_width_inches": round(label_width / 914400, 3),
        }
    )

    for index, row in enumerate(valid.itertuples(index=False)):
        y_mid = row_top + row_spacing * index + Inches(0.18)
        label = clean_text(getattr(row, "label"))
        add_fitted_textbox(
            slide,
            left + Inches(0.16),
            y_mid - label_height / 2,
            label_width - Inches(0.04),
            label_height,
            label,
            min_font_size=8,
            max_font_size=label_font_max,
            max_lines=1,
            allow_truncation=True,
            bold=row.kind == "selection",
            color=theme.text,
            font_name=theme.font_family_bold if row.kind == "selection" else theme.font_family,
        )
        track = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            plot_left,
            y_mid - Inches(0.015),
            plot_right - plot_left,
            Inches(0.03),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = theme.soft_line
        track.line.color.rgb = theme.soft_line
        if row.kind == "range":
            start = _scale_to_width(float(row.value_min), axis_min, axis_max, plot_left, plot_right)
            end = _scale_to_width(float(row.value_max), axis_min, axis_max, plot_left, plot_right)
            band = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                start,
                y_mid - Inches(0.05),
                max(Inches(0.03), end - start),
                Inches(0.1),
            )
            band.fill.solid()
            band.fill.fore_color.rgb = theme.teal
            band.fill.transparency = 0.15
            band.line.color.rgb = theme.teal
            for value in [float(row.value_min), float(row.value_max)]:
                marker_left = _scale_to_width(value, axis_min, axis_max, plot_left, plot_right) - Inches(0.04)
                marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, marker_left, y_mid - Inches(0.04), Inches(0.08), Inches(0.08))
                marker.fill.solid()
                marker.fill.fore_color.rgb = theme.teal
                marker.line.color.rgb = theme.surface
                marker.line.width = Pt(1)
            value_text = f"{format_indicator_value(row.value_min, unit)} to {format_indicator_value(row.value_max, unit)}"
        else:
            marker_center = _scale_to_width(float(row.value), axis_min, axis_max, plot_left, plot_right)
            value_line = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                plot_left,
                y_mid - Inches(0.02),
                max(Inches(0.03), marker_center - plot_left),
                Inches(0.04),
            )
            value_line.fill.solid()
            value_line.fill.fore_color.rgb = _kind_colour(row.kind, theme)
            value_line.line.color.rgb = _kind_colour(row.kind, theme)
            marker_left = marker_center - Inches(0.05)
            size = Inches(0.14) if row.kind == "selection" else Inches(0.11)
            marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, marker_left, y_mid - size / 2, size, size)
            marker.fill.solid()
            marker.fill.fore_color.rgb = _kind_colour(row.kind, theme)
            marker.line.color.rgb = theme.surface
            marker.line.width = Pt(1.2)
            value_text = format_indicator_value(row.value, unit)
        add_fitted_textbox(
            slide,
            plot_right + Inches(0.08),
            y_mid - value_height / 2,
            value_width - Inches(0.02),
            value_height,
            value_text,
            min_font_size=8,
            max_font_size=value_font_max,
            max_lines=1,
            allow_truncation=True,
            bold=row.kind == "selection",
            color=theme.text if row.kind == "selection" else theme.muted,
            align=1,
        )

    add_textbox(
        slide,
        plot_left,
        top + height - Inches(0.34),
        Inches(1.2),
        Inches(0.18),
        format_indicator_value(axis_min if unit not in {"share", "rate_per_1000"} else 0, unit),
        font_size=9,
        color=theme.muted,
    )
    add_textbox(
        slide,
        plot_right - Inches(1.1),
        top + height - Inches(0.34),
        Inches(1.2),
        Inches(0.18),
        format_indicator_value(axis_max, unit),
        font_size=9,
        color=theme.muted,
        align=1,
    )
    return debug


def render_distribution_band(
    slide,
    selection: dict[str, object] | None,
    benchmark_rows: pd.DataFrame,
    left,
    top,
    width,
    height,
    unit: str | None,
    *,
    theme: PptTheme,
) -> None:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.12),
        width - Inches(0.32),
        Inches(0.24),
        "Decile band",
        font_size=12,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )
    axis_left = left + Inches(0.6)
    axis_right = left + width - Inches(0.6)
    axis_y = top + Inches(1.2)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        axis_left,
        axis_y,
        axis_right - axis_left,
        Inches(0.22),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = theme.surface_alt
    band.line.color.rgb = theme.soft_line
    for decile in range(1, 11):
        x = _scale_to_width(float(decile), 1.0, 10.0, axis_left, axis_right)
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, axis_y - Inches(0.06), Inches(0.01), Inches(0.34))
        tick.fill.solid()
        tick.fill.fore_color.rgb = theme.soft_line
        tick.line.color.rgb = theme.soft_line
        add_textbox(
            slide,
            x - Inches(0.08),
            axis_y + Inches(0.25),
            Inches(0.16),
            Inches(0.18),
            str(decile),
            font_size=9,
            color=theme.muted,
            align=1,
        )
    add_textbox(slide, axis_left, axis_y - Inches(0.3), Inches(1.3), Inches(0.18), "1 most deprived", font_size=9, color=theme.muted)
    add_textbox(
        slide,
        axis_right - Inches(1.3),
        axis_y - Inches(0.3),
        Inches(1.3),
        Inches(0.18),
        "10 least deprived",
        font_size=9,
        color=theme.muted,
        align=1,
    )

    if selection and pd.notna(selection.get("value")):
        _add_band_marker(slide, axis_left, axis_right, axis_y, float(selection["value"]), "Selected footprint", theme.navy, theme)
    for row in benchmark_rows.itertuples(index=False):
        if row.kind == "borough":
            _add_band_marker(slide, axis_left, axis_right, axis_y + Inches(0.16), float(row.value), str(row.label), theme.teal, theme)
        elif row.kind == "london":
            _add_band_marker(slide, axis_left, axis_right, axis_y + Inches(0.34), float(row.value), "London", theme.orange, theme)

    if selection and pd.notna(selection.get("summary_min")) and pd.notna(selection.get("summary_max")):
        summary = f"Selected range: {format_indicator_value(selection['summary_min'], unit)} to {format_indicator_value(selection['summary_max'], unit)}"
        add_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(2.05),
            width - Inches(0.32),
            Inches(0.22),
            summary,
            font_size=11,
            color=theme.text,
        )


def render_trend_chart(slide, frame: pd.DataFrame, unit: str | None, left, top, width, height, *, theme: PptTheme) -> dict[str, object]:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.12),
        width - Inches(0.32),
        Inches(0.22),
        "Trend over time",
        font_size=12,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )

    chart_data = CategoryChartData()
    ordered_periods = [friendly_period(period) for period in frame["period"].drop_duplicates().tolist()]
    categories = _compact_month_labels(ordered_periods)
    chart_data.categories = categories

    if "series_kind" in frame.columns:
        selection_series = frame.loc[frame["series_kind"].astype(str) == "selection", "series"].drop_duplicates().astype(str).tolist()
        london_series = frame.loc[frame["series_kind"].astype(str) == "london", "series"].drop_duplicates().astype(str).tolist()
        comparison_series = frame.loc[
            ~frame["series_kind"].astype(str).isin({"selection", "london"}),
            "series",
        ].drop_duplicates().astype(str).tolist()
        series_order = selection_series[:1] + london_series[:1]
        additional = [name for name in comparison_series if name not in series_order]
    else:
        known_series = frame["series"].drop_duplicates().astype(str).tolist()
        primary_series = next((name for name in known_series if name != "London"), known_series[0] if known_series else "")
        series_order = [primary_series] if primary_series else []
        if "London" in known_series:
            series_order.append("London")
        additional = [name for name in known_series if name not in series_order]
    series_order.extend(additional[:2])
    max_series_label_chars = max((len(clean_text(name)) for name in series_order), default=0)
    compact_legend = len(series_order) >= theme.layout.trend_dense_series_threshold or max_series_label_chars > theme.layout.trend_compact_legend_chars
    categories_compacted = any(not label for label in categories)
    series_label_limit = 18 if compact_legend else 26

    for series_name in series_order:
        subset = frame[frame["series"] == series_name].copy()
        if subset.empty:
            continue
        series_values = [scale_for_chart(value, unit) or 0.0 for value in subset["value"]]
        chart_data.add_series(_compact_series_label(series_name, max_chars=series_label_limit), series_values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        left + Inches(0.18),
        top + Inches(0.42),
        width - Inches(0.36),
        height - Inches(0.6),
        chart_data,
    ).chart
    visible_series_count = len(chart.series)
    legend_position = XL_LEGEND_POSITION.BOTTOM if compact_legend else XL_LEGEND_POSITION.TOP
    chart.has_legend = visible_series_count > 1
    if chart.has_legend:
        chart.legend.position = legend_position
        chart.legend.include_in_layout = compact_legend
    chart.chart_style = 2
    chart.category_axis.tick_labels.font.name = theme.font_family
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.name = theme.font_family
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = chart_number_format(unit)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = theme.soft_line
    chart.category_axis.has_title = False
    chart.value_axis.has_title = False

    colours = [theme.navy, theme.orange, theme.teal, theme.blue, theme.gold]
    for idx, series in enumerate(chart.series):
        colour = colours[min(idx, len(colours) - 1)]
        series.format.line.color.rgb = colour
        series.format.line.width = Pt(2.6 if idx == 0 else 1.8)
        series.marker.size = 5 if idx == 0 else 4
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = colour
        series.marker.format.line.color.rgb = theme.surface
    return {
        "visible_series_count": visible_series_count,
        "max_series_label_chars": max_series_label_chars,
        "legend_position": "bottom" if compact_legend else "top",
        "legend_compacted": compact_legend,
        "series_label_limit": series_label_limit,
        "categories_compacted": categories_compacted,
    }


def render_choropleth_map(
    slide,
    map_values: pd.DataFrame,
    current_ids: list[str],
    left,
    top,
    width,
    height,
    unit: str | None,
    map_style: str = "real",
    *,
    theme: PptTheme,
) -> VisualPlacement:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.12),
        width - Inches(0.32),
        Inches(0.22),
        "Hex map" if map_style == "hex" else "Neighbourhood map",
        font_size=12,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )
    values = pd.to_numeric(map_values["value"], errors="coerce").dropna()
    image_stream, source_width_px, source_height_px = _choropleth_image_stream(
        map_values=map_values,
        current_ids=current_ids,
        map_style=map_style,
        width_px=1280,
        height_px=720,
    )
    placement = add_visual_contained(
        slide,
        image_stream,
        left + Inches(0.12),
        top + Inches(0.4),
        width - Inches(0.24),
        height - Inches(0.74),
        source_size=(source_width_px, source_height_px),
    )
    if not values.empty:
        value_note = (
            f"{format_indicator_value(float(values.min()), unit)} to "
            f"{format_indicator_value(float(values.max()), unit)} across London"
        )
    else:
        value_note = "No mapped values available"
    selection_descriptor = resolve_footprint_label(current_ids)
    selection_verb = "are" if selection_descriptor.is_plural else "is"
    add_textbox(
        slide,
        left + Inches(0.16),
        top + height - Inches(0.24),
        width - Inches(0.32),
        Inches(0.18),
        f"{selection_descriptor.label} {selection_verb} outlined. {value_note}.",
        font_size=9,
        color=theme.muted,
    )
    return placement


def _choropleth_image_stream(
    *,
    map_values: pd.DataFrame,
    current_ids: list[str],
    map_style: str,
    width_px: int,
    height_px: int,
) -> tuple[BytesIO, int, int]:
    geography = load_map_geography("hex" if map_style == "hex" else "real")[["neighbourhood_id", "geometry"]].copy()
    overlay = load_hex_icb_geography()[["geometry"]].copy() if map_style == "hex" else None
    frame = geography.merge(map_values[["neighbourhood_id", "value"]], on="neighbourhood_id", how="left")
    frame["value"] = pd.to_numeric(frame["value"], errors="coerce")

    image = Image.new("RGB", (width_px, height_px), (255, 255, 255))
    draw = ImageDraw.Draw(image)
    minx, miny, maxx, maxy = frame.total_bounds
    pad = 18
    scale_x = (width_px - 2 * pad) / max(maxx - minx, 1)
    scale_y = (height_px - 2 * pad) / max(maxy - miny, 1)
    scale = min(scale_x, scale_y)
    x_offset = (width_px - (maxx - minx) * scale) / 2
    y_offset = (height_px - (maxy - miny) * scale) / 2

    def transform(x: float, y: float) -> tuple[float, float]:
        return (
            x_offset + (x - minx) * scale,
            y_offset + (maxy - y) * scale,
        )

    values = frame["value"].dropna()
    value_min = float(values.min()) if not values.empty else 0.0
    value_max = float(values.max()) if not values.empty else 1.0
    selected_set = {str(value) for value in current_ids}

    for row in frame.itertuples(index=False):
        fill = _map_fill_colour(getattr(row, "value"), value_min, value_max)
        outline = (200, 107, 42) if str(row.neighbourhood_id) in selected_set else (255, 255, 255)
        line_width = 4 if str(row.neighbourhood_id) in selected_set else 1
        _draw_geometry(draw, row.geometry, transform, fill, outline, line_width)
    if overlay is not None and not overlay.empty:
        for row in overlay.itertuples(index=False):
            _draw_geometry_outline(draw, row.geometry, transform, (42, 18, 69), 3)

    diff = ImageChops.difference(image, Image.new("RGB", image.size, (255, 255, 255)))
    bbox = diff.getbbox()
    if bbox is not None:
        pad = 24
        left_crop = max(0, bbox[0] - pad)
        top_crop = max(0, bbox[1] - pad)
        right_crop = min(image.width, bbox[2] + pad)
        bottom_crop = min(image.height, bbox[3] + pad)
        image = image.crop((left_crop, top_crop, right_crop, bottom_crop))

    stream = BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    return stream, image.width, image.height


def _draw_geometry(draw: ImageDraw.ImageDraw, geometry, transform, fill, outline, line_width: int) -> None:
    if geometry is None or geometry.is_empty:
        return
    if geometry.geom_type == "Polygon":
        points = [transform(x, y) for x, y in geometry.exterior.coords]
        draw.polygon(points, fill=fill)
        draw.line(points + [points[0]], fill=outline, width=line_width)
        return
    if geometry.geom_type == "MultiPolygon":
        for part in geometry.geoms:
            _draw_geometry(draw, part, transform, fill, outline, line_width)


def _draw_geometry_outline(draw: ImageDraw.ImageDraw, geometry, transform, outline, line_width: int) -> None:
    if geometry is None or geometry.is_empty:
        return
    if geometry.geom_type == "Polygon":
        points = [transform(x, y) for x, y in geometry.exterior.coords]
        draw.line(points + [points[0]], fill=outline, width=line_width)
        return
    if geometry.geom_type == "MultiPolygon":
        for part in geometry.geoms:
            _draw_geometry_outline(draw, part, transform, outline, line_width)


def _map_fill_colour(value: float | int | None, value_min: float, value_max: float) -> tuple[int, int, int]:
    if value is None or pd.isna(value):
        return (229, 238, 244)
    if value_max <= value_min:
        ratio = 0.65
    else:
        ratio = (float(value) - value_min) / (value_max - value_min)
    low = (229, 238, 244)
    high = (15, 59, 95)
    return tuple(int(round(low[idx] + (high[idx] - low[idx]) * ratio)) for idx in range(3))


def render_context_strip(slide, distribution: pd.DataFrame, value: float | int | None, unit: str | None, left, top, width, height, *, theme: PptTheme) -> None:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(0.1),
        width - Inches(0.28),
        Inches(0.2),
        "London context",
        font_size=11,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )
    if distribution.empty or value is None or pd.isna(value):
        add_textbox(
            slide,
            left + Inches(0.14),
            top + Inches(0.44),
            width - Inches(0.28),
            Inches(0.24),
            "No London-wide distribution is available for this indicator.",
            font_size=10,
            color=theme.muted,
        )
        return
    values = pd.to_numeric(distribution["value"], errors="coerce").dropna()
    if values.empty:
        return
    min_value = float(values.min())
    max_value = float(values.max())
    median_value = float(values.median())
    axis_left = left + Inches(0.24)
    axis_right = left + width - Inches(0.24)
    axis_y = top + Inches(0.84)
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, axis_left, axis_y, axis_right - axis_left, Inches(0.02))
    axis.fill.solid()
    axis.fill.fore_color.rgb = theme.soft_line
    axis.line.color.rgb = theme.soft_line
    for marker_value, colour, label in [
        (median_value, theme.gold, "Median"),
        (float(value), theme.navy, "Selected"),
    ]:
        x = _scale_to_width(marker_value, min_value, max_value, axis_left, axis_right)
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Inches(0.05), axis_y - Inches(0.05), Inches(0.1), Inches(0.1))
        marker.fill.solid()
        marker.fill.fore_color.rgb = colour
        marker.line.color.rgb = theme.surface
        add_textbox(
            slide,
            x - Inches(0.4),
            axis_y + Inches(0.08),
            Inches(0.8),
            Inches(0.18),
            label,
            font_size=8,
            color=theme.muted,
            align=1,
        )
    add_textbox(slide, axis_left, top + height - Inches(0.35), Inches(1.0), Inches(0.18), format_indicator_value(min_value, unit), font_size=8, color=theme.muted)
    add_textbox(slide, axis_right - Inches(1.0), top + height - Inches(0.35), Inches(1.0), Inches(0.18), format_indicator_value(max_value, unit), font_size=8, color=theme.muted, align=1)
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(1.12),
        width - Inches(0.28),
        Inches(0.42),
        percentile_note(distribution, value),
        font_size=10,
        color=theme.text,
    )


def render_population_pyramid(
    slide,
    composition_df: pd.DataFrame,
    *,
    comparator_subject: str,
    left,
    top,
    width,
    height,
    theme: PptTheme,
) -> dict[str, object]:
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.12),
        width - Inches(0.32),
        Inches(0.22),
        "Population pyramid",
        font_size=12,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )

    selection_subject, comparator_subjects = _population_pyramid_subject_options(composition_df)
    chosen_comparator = comparator_subject if comparator_subject in comparator_subjects else (comparator_subjects[0] if comparator_subjects else "")
    if not selection_subject or not chosen_comparator:
        add_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(0.52),
            width - Inches(0.32),
            Inches(0.3),
            "No comparator is available for the population pyramid.",
            font_size=11,
            color=theme.muted,
        )
        return {"rendered": False, "selection_subject": selection_subject, "comparator_subject": chosen_comparator}

    chart = _population_pyramid_frame(composition_df, selection_subject, chosen_comparator)
    if chart.empty:
        add_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(0.52),
            width - Inches(0.32),
            Inches(0.3),
            "No composition data are available for this view.",
            font_size=11,
            color=theme.muted,
        )
        return {"rendered": False, "selection_subject": selection_subject, "comparator_subject": chosen_comparator}

    max_share = max(
        float(chart["selected_share_pct"].max()),
        float(chart["comparator_share_pct"].max()),
        1.0,
    )
    axis_limit = max(5, int(((max_share + 4.999) // 5) * 5))
    center_x = left + width / 2
    chart_top = top + Inches(0.68)
    chart_bottom = top + height - Inches(0.52)
    row_spacing = min(Inches(0.34), (chart_bottom - chart_top) / max(len(chart), 1))
    bar_height = max(Inches(0.12), row_spacing - Inches(0.08))
    label_width = Inches(1.55)
    inner_gap = Inches(0.08)
    half_plot_width = max((width - label_width - Inches(0.9)) / 2, Inches(1.35))
    center_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        center_x - Inches(0.006),
        chart_top - Inches(0.06),
        Inches(0.012),
        chart_bottom - chart_top + Inches(0.1),
    )
    center_line.fill.solid()
    center_line.fill.fore_color.rgb = theme.soft_line
    center_line.line.color.rgb = theme.soft_line

    add_textbox(
        slide,
        center_x - label_width - half_plot_width,
        top + Inches(0.38),
        half_plot_width,
        Inches(0.18),
        selection_subject,
        font_size=10,
        color=theme.navy,
        align=1,
    )
    add_textbox(
        slide,
        center_x + inner_gap,
        top + Inches(0.38),
        half_plot_width,
        Inches(0.18),
        chosen_comparator,
        font_size=10,
        color=theme.teal,
        align=1,
    )

    for index, row in enumerate(chart.itertuples(index=False)):
        y_top = chart_top + row_spacing * index + (row_spacing - bar_height) / 2
        y_mid = y_top + bar_height / 2
        add_fitted_textbox(
            slide,
            center_x - label_width / 2,
            y_top - Inches(0.01),
            label_width,
            bar_height + Inches(0.02),
            str(row.category),
            min_font_size=7,
            max_font_size=9,
            max_lines=1,
            allow_truncation=True,
            align=1,
            color=theme.text,
        )

        selected_width = half_plot_width * (float(row.selected_share_pct) / axis_limit)
        comparator_width = half_plot_width * (float(row.comparator_share_pct) / axis_limit)

        selected_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            center_x - inner_gap - selected_width,
            y_top,
            max(selected_width, Inches(0.02)),
            bar_height,
        )
        selected_bar.fill.solid()
        selected_bar.fill.fore_color.rgb = theme.navy
        selected_bar.line.color.rgb = theme.navy

        comparator_bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            center_x + inner_gap,
            y_top,
            max(comparator_width, Inches(0.02)),
            bar_height,
        )
        comparator_bar.fill.solid()
        comparator_bar.fill.fore_color.rgb = theme.teal
        comparator_bar.line.color.rgb = theme.teal

        add_fitted_textbox(
            slide,
            center_x - inner_gap - selected_width - Inches(0.56),
            y_mid - Inches(0.08),
            Inches(0.5),
            Inches(0.16),
            f"{float(row.selected_share_pct):.1f}%",
            min_font_size=7,
            max_font_size=8,
            max_lines=1,
            allow_truncation=False,
            align=2,
            color=theme.muted,
        )
        add_fitted_textbox(
            slide,
            center_x + inner_gap + comparator_width + Inches(0.04),
            y_mid - Inches(0.08),
            Inches(0.5),
            Inches(0.16),
            f"{float(row.comparator_share_pct):.1f}%",
            min_font_size=7,
            max_font_size=8,
            max_lines=1,
            allow_truncation=False,
            color=theme.muted,
        )

    add_textbox(slide, center_x - label_width - half_plot_width, chart_bottom + Inches(0.04), Inches(0.35), Inches(0.16), f"{axis_limit}%", font_size=8, color=theme.muted)
    add_textbox(slide, center_x - Inches(0.12), chart_bottom + Inches(0.04), Inches(0.24), Inches(0.16), "0%", font_size=8, color=theme.muted, align=1)
    add_textbox(slide, center_x + half_plot_width + Inches(0.12), chart_bottom + Inches(0.04), Inches(0.35), Inches(0.16), f"{axis_limit}%", font_size=8, color=theme.muted, align=2)

    return {
        "rendered": True,
        "selection_subject": selection_subject,
        "comparator_subject": chosen_comparator,
        "category_count": int(len(chart)),
    }


def _population_pyramid_subject_options(composition_df: pd.DataFrame) -> tuple[str, list[str]]:
    if composition_df.empty:
        return "", []
    ordered = composition_df.sort_values("subject_order")[["subject", "subject_kind"]].drop_duplicates()
    subject_order = ordered["subject"].astype(str).tolist()
    selection_subject = next(
        (str(row.subject) for row in ordered.itertuples(index=False) if str(row.subject_kind) == "selection"),
        subject_order[0] if subject_order else "",
    )
    comparator_subjects = [subject for subject in subject_order if subject and subject != selection_subject]
    return selection_subject, comparator_subjects


def _population_pyramid_frame(composition_df: pd.DataFrame, selection_subject: str, comparator_subject: str) -> pd.DataFrame:
    frame = composition_df.copy()
    frame["share_pct"] = pd.to_numeric(frame["share"], errors="coerce").fillna(0.0) * 100.0
    selected_df = (
        frame[frame["subject"].astype(str) == selection_subject][["category", "category_order", "share_pct"]]
        .rename(columns={"share_pct": "selected_share_pct"})
        .copy()
    )
    comparator_df = (
        frame[frame["subject"].astype(str) == comparator_subject][["category", "category_order", "share_pct"]]
        .rename(columns={"share_pct": "comparator_share_pct"})
        .copy()
    )
    return (
        selected_df.merge(comparator_df, on=["category", "category_order"], how="outer")
        .fillna({"selected_share_pct": 0.0, "comparator_share_pct": 0.0})
        .sort_values("category_order", ascending=False)
        .reset_index(drop=True)
    )


def _kind_colour(kind: str, theme: PptTheme):
    if kind == "selection":
        return theme.navy
    if kind == "borough":
        return theme.teal
    if kind == "london":
        return theme.orange
    return theme.blue


def _scale_to_width(value: float, min_value: float, max_value: float, left, right):
    span = max_value - min_value if max_value != min_value else 1.0
    return left + (right - left) * ((value - min_value) / span)


def _add_band_marker(slide, axis_left, axis_right, axis_y, value: float, label: str, colour, theme: PptTheme) -> None:
    x = _scale_to_width(value, 1.0, 10.0, axis_left, axis_right)
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Inches(0.05), axis_y - Inches(0.05), Inches(0.1), Inches(0.1))
    marker.fill.solid()
    marker.fill.fore_color.rgb = colour
    marker.line.color.rgb = theme.surface
    add_textbox(
        slide,
        x - Inches(0.45),
        axis_y - Inches(0.28),
        Inches(0.9),
        Inches(0.16),
        label,
        font_size=8,
        color=theme.muted,
        align=1,
    )


def _compact_month_labels(labels: list[str]) -> list[str]:
    if len(labels) <= 12:
        return labels
    compact: list[str] = []
    for idx, label in enumerate(labels):
        if idx == 0 or idx == len(labels) - 1 or idx % 3 == 0:
            compact.append(label)
        else:
            compact.append("")
    return compact


def _compact_series_label(name: object, *, max_chars: int) -> str:
    value = clean_text(name)
    replacements = {
        "London Borough of ": "",
        "Neighbourhoods": "Nhds",
        "neighbourhoods": "nhds",
        "benchmark": "bench.",
        "benchmarks": "bench.",
    }
    for source_text, target_text in replacements.items():
        value = value.replace(source_text, target_text)
    return truncate_text(value, max_chars)
