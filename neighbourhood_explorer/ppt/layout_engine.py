from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
import math
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from neighbourhood_explorer.paths import project_path

from .indicator_layouts import Box
from .style_tokens import PptTheme, THEME


LOGO_PATH = project_path("assets", "logos", "PPL Logo_RGB.png")


EMU_PER_INCH = 914400


@dataclass
class TextFitResult:
    text: str
    font_size_pt: int
    line_spacing: float
    line_count: int
    truncated: bool
    overflow: bool
    estimated_height: int
    box: Box


@dataclass
class VisualPlacement:
    box: Box
    source_width_px: int
    source_height_px: int
    rendered_width: int
    rendered_height: int
    distorted: bool = False


def create_presentation(theme: PptTheme = THEME) -> Presentation:
    prs = Presentation()
    prs.slide_width = theme.slide_width
    prs.slide_height = theme.slide_height
    return prs


def apply_slide_background(slide, theme: PptTheme = THEME) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = theme.background


def add_logo(slide, theme: PptTheme = THEME, *, right_aligned: bool = False, width: int | None = None) -> None:
    logo = Path(LOGO_PATH)
    if not logo.exists():
        return
    display_width = width or Inches(1.1)
    left = theme.margin_x if not right_aligned else theme.slide_width - theme.margin_x - display_width
    slide.shapes.add_picture(str(logo), left, theme.margin_y - Inches(0.03), width=display_width)


def add_panel(
    slide,
    left,
    top,
    width,
    height,
    *,
    theme: PptTheme = THEME,
    fill=None,
    line=None,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or theme.surface
    shape.line.color.rgb = line or theme.soft_line
    shape.line.width = Pt(1)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size: int = 18,
    bold: bool = False,
    color=None,
    font_name: str | None = None,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    margin: float = 0.03,
    theme: PptTheme = THEME,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font_name or theme.font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or theme.text
    return box


def estimate_text_height(
    text: object,
    width: int,
    font_size: int,
    *,
    max_lines: int | None = None,
    margin: float = 0.03,
    line_spacing: float = 1.08,
) -> int:
    line_count = _estimate_wrapped_line_count(text, width, font_size, margin=margin, max_lines=max_lines)
    return int(Inches(margin * 2) + Pt(font_size * line_spacing * max(line_count, 1)))


def truncate_to_lines(
    text: object,
    width: int,
    font_size: int,
    *,
    max_lines: int,
    margin: float = 0.03,
) -> tuple[str, int, bool]:
    wrapped_lines = _wrap_text_lines(text, width, font_size, margin=margin)
    if len(wrapped_lines) <= max_lines:
        value = "\n".join(wrapped_lines)
        return value, len(wrapped_lines), False
    clipped = wrapped_lines[:max_lines]
    clipped[-1] = clipped[-1].rstrip(" ,;:.") + "..."
    return "\n".join(clipped), max_lines, True


def choose_font_size_for_box(
    text: object,
    width: int,
    height: int,
    *,
    min_font_size: int,
    max_font_size: int,
    max_lines: int,
    allow_truncation: bool = True,
    margin: float = 0.03,
    line_spacing: float = 1.08,
) -> TextFitResult:
    box = Box(0, 0, width, height)
    value = str(text or "")
    spacing_options = []
    for candidate in [line_spacing, max(0.98, line_spacing - 0.08), max(0.92, line_spacing - 0.14)]:
        rounded = round(candidate, 2)
        if rounded not in spacing_options:
            spacing_options.append(rounded)
    fallback = TextFitResult(value, min_font_size, spacing_options[0], 1, False, False, 0, box)
    for font_size in range(max_font_size, min_font_size - 1, -1):
        rendered_text, line_count, truncated = (
            truncate_to_lines(
                value,
                width,
                font_size,
                max_lines=max_lines,
                margin=margin,
            )
            if allow_truncation
            else (
                value,
                _estimate_wrapped_line_count(value, width, font_size, margin=margin, max_lines=max_lines),
                False,
            )
        )
        for spacing in spacing_options:
            estimated_height = estimate_text_height(
                rendered_text,
                width,
                font_size,
                max_lines=max_lines,
                margin=margin,
                line_spacing=spacing,
            )
            overflow = estimated_height > height
            candidate = TextFitResult(rendered_text, font_size, spacing, line_count, truncated, overflow, estimated_height, box)
            fallback = candidate
            if not overflow:
                return candidate
    return fallback


def fit_text_to_box(
    text: object,
    box: Box,
    *,
    min_font_size: int,
    max_font_size: int,
    max_lines: int,
    allow_truncation: bool = True,
    margin: float = 0.03,
    line_spacing: float = 1.08,
) -> TextFitResult:
    result = choose_font_size_for_box(
        text,
        box.width,
        box.height,
        min_font_size=min_font_size,
        max_font_size=max_font_size,
        max_lines=max_lines,
        allow_truncation=allow_truncation,
        margin=margin,
        line_spacing=line_spacing,
    )
    result.box = box
    return result


def add_fitted_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    min_font_size: int,
    max_font_size: int,
    max_lines: int,
    allow_truncation: bool = True,
    bold: bool = False,
    color=None,
    font_name: str | None = None,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    margin: float = 0.03,
    line_spacing: float = 1.08,
    theme: PptTheme = THEME,
):
    box = Box(left, top, width, height)
    fit = fit_text_to_box(
        text,
        box,
        min_font_size=min_font_size,
        max_font_size=max_font_size,
        max_lines=max_lines,
        allow_truncation=allow_truncation,
        margin=margin,
        line_spacing=line_spacing,
    )
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = Pt(fit.font_size_pt * fit.line_spacing)
    run = paragraph.add_run()
    run.text = fit.text
    run.font.name = font_name or theme.font_family
    run.font.size = Pt(fit.font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = color or theme.text
    return shape, fit


def add_visual_contained(
    slide,
    image_source,
    left,
    top,
    width,
    height,
    *,
    source_size: tuple[int, int] | None = None,
) -> VisualPlacement:
    image_stream, source_width_px, source_height_px = _normalise_image_source(image_source, source_size=source_size)
    source_ratio = source_width_px / max(source_height_px, 1)
    target_ratio = width / max(height, 1)
    scale = min(width / max(source_width_px, 1), height / max(source_height_px, 1))
    rendered_width = int(source_width_px * scale)
    rendered_height = int(source_height_px * scale)
    render_left = int(left + max(0, (width - rendered_width) / 2))
    render_top = int(top + max(0, (height - rendered_height) / 2))
    slide.shapes.add_picture(image_stream, render_left, render_top, width=rendered_width, height=rendered_height)
    rendered_ratio = rendered_width / max(rendered_height, 1)
    distorted = not math.isclose(source_ratio, rendered_ratio, rel_tol=0.01, abs_tol=0.01)
    return VisualPlacement(
        box=Box(render_left, render_top, rendered_width, rendered_height),
        source_width_px=source_width_px,
        source_height_px=source_height_px,
        rendered_width=rendered_width,
        rendered_height=rendered_height,
        distorted=distorted,
    )


def build_layout_preview(
    plan,
    *,
    slide_id: str = "",
    validation=None,
    text_results: dict[str, TextFitResult] | None = None,
    theme: PptTheme = THEME,
    pixels_per_inch: int = 72,
) -> bytes:
    width_px = max(1, int((theme.slide_width / EMU_PER_INCH) * pixels_per_inch))
    height_px = max(1, int((theme.slide_height / EMU_PER_INCH) * pixels_per_inch))
    image = Image.new("RGB", (width_px, height_px), _rgb_tuple(theme.background))
    draw = ImageDraw.Draw(image)

    truncated_flags, overflow_flags = _collect_box_flags(plan, text_results or {})
    box_fills = {
        "title": _rgb_tuple(theme.surface_tint),
        "subtitle": _rgb_tuple(theme.surface_alt),
        "visual": _rgb_tuple(theme.surface),
        "footer_note": _rgb_tuple(theme.surface_alt),
        "footer_meta": _rgb_tuple(theme.surface_alt),
        "footer_band": _rgb_tuple(theme.surface_tint),
        "primary_narrative": _rgb_tuple(theme.surface),
        "secondary_narrative": _rgb_tuple(theme.surface),
        "merged_narrative": _rgb_tuple(theme.surface),
    }
    for name, box in plan.occupied_boxes():
        rect = _scale_box_to_preview(box, pixels_per_inch)
        fill = box_fills.get(name, _rgb_tuple(theme.surface))
        outline = _rgb_tuple(theme.soft_line)
        if name in truncated_flags:
            outline = (214, 137, 16)
        if name in overflow_flags:
            outline = (212, 53, 28)
        draw.rounded_rectangle(rect, radius=8, fill=fill, outline=outline, width=2)
        draw.text((rect[0] + 6, rect[1] + 5), name.replace("_", " "), fill=_rgb_tuple(theme.text))

    if slide_id:
        draw.text((10, 8), slide_id, fill=_rgb_tuple(theme.text))
    if validation is not None:
        status = f"Warnings: {len(getattr(validation, 'warnings', []))} | Fullness: {getattr(validation, 'fullness_score', 0):.2f}"
        draw.text((10, height_px - 36), status, fill=_rgb_tuple(theme.muted))
        warning_text = " | ".join(getattr(validation, "warnings", [])[:2])
        if warning_text:
            draw.text((10, height_px - 20), warning_text, fill=(212, 53, 28))

    stream = BytesIO()
    image.save(stream, format="PNG")
    return stream.getvalue()


def add_badge(slide, left, top, text, *, fill, color, theme: PptTheme = THEME, width: int | None = None):
    badge_width = width or Inches(max(1.2, min(4.0, 0.12 * max(len(str(text)), 8))))
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, badge_width, Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    badge.line.color.rgb = fill
    frame = badge.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = theme.font_family
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    return badge


def _scale_box_to_preview(box: Box, pixels_per_inch: int) -> tuple[int, int, int, int]:
    scale = pixels_per_inch / EMU_PER_INCH
    return (
        int(box.left * scale),
        int(box.top * scale),
        int(box.right * scale),
        int(box.bottom * scale),
    )


def _collect_box_flags(plan, text_results: dict[str, TextFitResult]) -> tuple[set[str], set[str]]:
    truncated: set[str] = set()
    overflow: set[str] = set()
    box_names = [name for name, _ in plan.occupied_boxes()]
    for key, fit in text_results.items():
        matched_name = next((name for name in box_names if key == name or key.startswith(f"{name}_")), None)
        if matched_name is None:
            continue
        if getattr(fit, "truncated", False):
            truncated.add(matched_name)
        if getattr(fit, "overflow", False):
            overflow.add(matched_name)
    return truncated, overflow


def _rgb_tuple(colour) -> tuple[int, int, int]:
    return tuple(int(channel) for channel in colour)


def add_metric_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    label: str,
    value: str,
    caption: str = "",
    accent=None,
    theme: PptTheme = THEME,
) -> dict[str, TextFitResult | None]:
    accent_colour = accent or theme.blue
    add_panel(slide, left, top, width, height, theme=theme, fill=theme.surface, line=theme.soft_line)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_colour
    stripe.line.color.rgb = accent_colour
    _, label_fit = add_fitted_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.08),
        width - Inches(0.22),
        Inches(0.24),
        label,
        min_font_size=8,
        max_font_size=11,
        max_lines=2,
        allow_truncation=True,
        bold=True,
        color=theme.muted,
        font_name=theme.font_family_bold,
    )
    _, value_fit = add_fitted_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.33),
        width - Inches(0.22),
        Inches(0.35),
        value,
        min_font_size=14,
        max_font_size=22 if len(value) < 16 else 18,
        max_lines=2,
        allow_truncation=False,
        bold=True,
        color=theme.text,
        font_name=theme.font_family_bold,
    )
    caption_fit: TextFitResult | None = None
    if caption:
        _, caption_fit = add_fitted_textbox(
            slide,
            left + Inches(0.16),
            top + Inches(0.72),
            width - Inches(0.22),
            Inches(0.22),
            caption,
            min_font_size=7,
            max_font_size=9,
            max_lines=2,
            allow_truncation=True,
            color=theme.muted,
        )
    return {
        "label": label_fit,
        "value": value_fit,
        "caption": caption_fit,
    }


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    font_size: int = 12,
    theme: PptTheme = THEME,
    color=None,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(0)
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = f"• {item}"
        run.font.name = theme.font_family
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or theme.text


def add_page_number(slide, page_number: int, theme: PptTheme = THEME) -> None:
    add_textbox(
        slide,
        theme.slide_width - Inches(0.65),
        theme.slide_height - Inches(0.42),
        Inches(0.28),
        Inches(0.16),
        str(page_number),
        font_size=9,
        color=theme.muted,
        align=PP_ALIGN.RIGHT,
    )


def add_footer_rule(slide, theme: PptTheme = THEME) -> None:
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        theme.margin_x,
        theme.slide_height - Inches(0.78),
        theme.slide_width - 2 * theme.margin_x,
        Inches(0.01),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.soft_line
    rule.line.color.rgb = theme.soft_line


def _estimate_wrapped_line_count(
    text: object,
    width: int,
    font_size: int,
    *,
    margin: float = 0.03,
    max_lines: int | None = None,
) -> int:
    lines = _wrap_text_lines(text, width, font_size, margin=margin)
    if max_lines is None:
        return max(1, len(lines))
    return min(max_lines, max(1, len(lines)))


def _wrap_text_lines(text: object, width: int, font_size: int, *, margin: float = 0.03) -> list[str]:
    value = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
    paragraphs = [paragraph.strip() for paragraph in value.split("\n") if paragraph.strip()] or [""]
    wrapped_lines: list[str] = []
    chars_per_line = max(8, _estimate_chars_per_line(width, font_size, margin=margin))
    for paragraph in paragraphs:
        words = paragraph.split()
        if not words:
            wrapped_lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}".strip()
            if len(candidate) <= chars_per_line:
                current = candidate
            else:
                wrapped_lines.append(current)
                current = word
        wrapped_lines.append(current)
    return wrapped_lines or [""]


def _estimate_chars_per_line(width: int, font_size: int, *, margin: float = 0.03) -> int:
    usable_width_inches = max(0.5, (width / EMU_PER_INCH) - (margin * 2))
    average_character_width = max(0.055, font_size / 170)
    return int(usable_width_inches / average_character_width)


def _normalise_image_source(image_source, *, source_size: tuple[int, int] | None = None) -> tuple[BytesIO, int, int]:
    if hasattr(image_source, "getvalue"):
        payload = image_source.getvalue()
    elif hasattr(image_source, "read"):
        payload = image_source.read()
    else:
        payload = Path(image_source).read_bytes()
    stream = BytesIO(payload)
    if source_size is None:
        with Image.open(BytesIO(payload)) as image:
            width_px, height_px = image.size
    else:
        width_px, height_px = source_size
    stream.seek(0)
    return stream, int(width_px), int(height_px)
